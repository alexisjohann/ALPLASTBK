#!/usr/bin/env python3
"""
FehrAdvice Slide Master Template Generator
==========================================

Erstellt ein PPTX Master-Template basierend auf dem Corporate Identity Guide
(Corporate-Identity-Guide-FAP-EN-231007 2.pdf, Seite 17)

Layouts aus dem CI Guide:
1. TITLE - Titelfolie mit Bild
2. SECTION - Abschnittsteiler
3. CONTENT - Standard-Inhalt mit Titel
4. DATA_CHART - Charts und Diagramme
5. KPI - Kennzahlen-Darstellung
6. TABLE - Tabellen
7. PROCESS - Prozess/Timeline
8. QUOTE - Zitat/Statement

Version: 1.0
Date: January 2026
"""

import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import nsmap
    from pptx.oxml import parse_xml
    from pptx.slide import SlideLayout, SlideMaster
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Error: python-pptx nicht installiert. Installiere mit: pip install python-pptx")
    sys.exit(1)

# =============================================================================
# FEHRADVICE CORPORATE COLORS (from CI Guide Page 9-10)
# =============================================================================

COLORS = {
    # Primary Colors
    "darkblue": RGBColor(0x02, 0x40, 0x79),      # #024079
    "lightblue": RGBColor(0x54, 0x9E, 0xDE),     # #549EDE
    "darkgray": RGBColor(0x25, 0x21, 0x2A),      # #25212A (text)
    "lightgray": RGBColor(0xF3, 0xF5, 0xF7),     # #F3F5F7 (background)
    "white": RGBColor(0xFF, 0xFF, 0xFF),

    # Secondary Colors
    "lilac": RGBColor(0xA1, 0xA0, 0xC6),         # #A1A0C6
    "mint": RGBColor(0x7E, 0xBD, 0xAC),          # #7EBDAC
    "ocher": RGBColor(0xDE, 0xCB, 0x3F),         # #DECB3F
    "orange": RGBColor(0xDE, 0x9D, 0x3E),        # #DE9D3E
}

# =============================================================================
# SLIDE DIMENSIONS (Widescreen 16:9)
# =============================================================================

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Logo paths
ASSETS_DIR = Path(__file__).parent.parent / "assets"
LOGO_WHITE = ASSETS_DIR / "logo" / "monochrome" / "fehradvice-logo-white.png"
LOGO_COLOR = ASSETS_DIR / "logo" / "color" / "fehradvice-logo-color-tagline.png"

# =============================================================================
# LAYOUT DEFINITIONS (from CI Guide Page 17)
# =============================================================================

LAYOUTS = {
    "TITLE": {
        "name": "FA Title Slide",
        "description": "Titelfolie mit Bild und Untertitel",
        "elements": {
            "background": {"type": "image", "position": "full"},
            "title": {"x": 0.5, "y": 2.0, "w": 6.0, "h": 1.5, "font_size": 44, "bold": True, "color": "white"},
            "subtitle": {"x": 0.5, "y": 3.5, "w": 6.0, "h": 0.8, "font_size": 20, "color": "white"},
            "logo": {"x": 11.5, "y": 0.3, "w": 1.5, "h": 0.6},
        }
    },
    "SECTION": {
        "name": "FA Section Divider",
        "description": "Abschnittsteiler mit hellblauem Hintergrund",
        "elements": {
            "background": {"type": "solid", "color": "lightblue"},
            "title": {"x": 2.0, "y": 2.5, "w": 9.333, "h": 1.5, "font_size": 48, "bold": True, "color": "darkblue", "align": "center"},
            "subtitle": {"x": 2.0, "y": 4.0, "w": 9.333, "h": 1.0, "font_size": 24, "color": "darkblue", "align": "center"},
            "logo": {"x": 11.5, "y": 0.3, "w": 1.5, "h": 0.6},
            "page_info": {"x": 0.3, "y": 0.3, "w": 4.0, "h": 0.4, "font_size": 10, "color": "darkblue"},
        }
    },
    "CONTENT": {
        "name": "FA Content",
        "description": "Standard-Inhaltsfolie mit Titel",
        "elements": {
            "header_line": {"x": 0, "y": 0.9, "w": 13.333, "h": 0.02, "color": "darkblue"},
            "title": {"x": 0.5, "y": 0.3, "w": 10.0, "h": 0.6, "font_size": 28, "bold": True, "color": "darkblue"},
            "content": {"x": 0.5, "y": 1.2, "w": 12.333, "h": 5.5, "font_size": 16, "color": "darkgray"},
            "logo": {"x": 11.5, "y": 0.15, "w": 1.5, "h": 0.6},
            "footer": {"x": 0.5, "y": 7.0, "w": 12.333, "h": 0.3, "font_size": 9, "color": "darkgray"},
        }
    },
    "DATA_CHART": {
        "name": "FA Data Chart",
        "description": "Folie mit Chart-Platzhalter",
        "elements": {
            "header_line": {"x": 0, "y": 0.9, "w": 13.333, "h": 0.02, "color": "darkblue"},
            "title": {"x": 0.5, "y": 0.3, "w": 10.0, "h": 0.6, "font_size": 28, "bold": True, "color": "darkblue"},
            "chart_placeholder": {"x": 0.5, "y": 1.2, "w": 12.333, "h": 5.5, "type": "picture"},
            "logo": {"x": 11.5, "y": 0.15, "w": 1.5, "h": 0.6},
            "footer": {"x": 0.5, "y": 7.0, "w": 12.333, "h": 0.3, "font_size": 9, "color": "darkgray"},
        }
    },
    "KPI": {
        "name": "FA KPI Dashboard",
        "description": "4 KPI-Boxen mit grossen Zahlen",
        "elements": {
            "header_line": {"x": 0, "y": 0.9, "w": 13.333, "h": 0.02, "color": "darkblue"},
            "title": {"x": 0.5, "y": 0.3, "w": 10.0, "h": 0.6, "font_size": 28, "bold": True, "color": "darkblue"},
            "kpi_1": {"x": 0.5, "y": 1.5, "w": 3.0, "h": 2.5, "type": "kpi_box"},
            "kpi_2": {"x": 3.666, "y": 1.5, "w": 3.0, "h": 2.5, "type": "kpi_box"},
            "kpi_3": {"x": 6.833, "y": 1.5, "w": 3.0, "h": 2.5, "type": "kpi_box"},
            "kpi_4": {"x": 10.0, "y": 1.5, "w": 3.0, "h": 2.5, "type": "kpi_box"},
            "logo": {"x": 11.5, "y": 0.15, "w": 1.5, "h": 0.6},
            "footer": {"x": 0.5, "y": 7.0, "w": 12.333, "h": 0.3, "font_size": 9, "color": "darkgray"},
        }
    },
    "TWO_COLUMN": {
        "name": "FA Two Column",
        "description": "Zwei-Spalten-Layout (Text + Bild/Chart)",
        "elements": {
            "header_line": {"x": 0, "y": 0.9, "w": 13.333, "h": 0.02, "color": "darkblue"},
            "title": {"x": 0.5, "y": 0.3, "w": 10.0, "h": 0.6, "font_size": 28, "bold": True, "color": "darkblue"},
            "left_content": {"x": 0.5, "y": 1.2, "w": 5.5, "h": 5.5, "font_size": 16, "color": "darkgray"},
            "right_content": {"x": 6.5, "y": 1.2, "w": 6.333, "h": 5.5, "type": "picture"},
            "logo": {"x": 11.5, "y": 0.15, "w": 1.5, "h": 0.6},
            "footer": {"x": 0.5, "y": 7.0, "w": 12.333, "h": 0.3, "font_size": 9, "color": "darkgray"},
        }
    },
    "PROCESS": {
        "name": "FA Process/Timeline",
        "description": "Prozess mit nummerierten Schritten",
        "elements": {
            "header_line": {"x": 0, "y": 0.9, "w": 13.333, "h": 0.02, "color": "darkblue"},
            "title": {"x": 0.5, "y": 0.3, "w": 10.0, "h": 0.6, "font_size": 28, "bold": True, "color": "darkblue"},
            "timeline_line": {"x": 1.0, "y": 4.0, "w": 11.333, "h": 0.05, "color": "lightblue"},
            "step_1": {"x": 1.5, "y": 2.0, "w": 2.0, "h": 3.5, "type": "step_box", "number": "01"},
            "step_2": {"x": 4.0, "y": 2.0, "w": 2.0, "h": 3.5, "type": "step_box", "number": "02"},
            "step_3": {"x": 6.5, "y": 2.0, "w": 2.0, "h": 3.5, "type": "step_box", "number": "03"},
            "step_4": {"x": 9.0, "y": 2.0, "w": 2.0, "h": 3.5, "type": "step_box", "number": "04"},
            "step_5": {"x": 11.5, "y": 2.0, "w": 1.5, "h": 3.5, "type": "step_box", "number": "05"},
            "logo": {"x": 11.5, "y": 0.15, "w": 1.5, "h": 0.6},
            "footer": {"x": 0.5, "y": 7.0, "w": 12.333, "h": 0.3, "font_size": 9, "color": "darkgray"},
        }
    },
    "QUOTE": {
        "name": "FA Quote/Statement",
        "description": "Zitat oder wichtige Aussage",
        "elements": {
            "background": {"type": "solid", "color": "lightblue"},
            "quote": {"x": 1.0, "y": 2.5, "w": 8.0, "h": 3.0, "font_size": 32, "bold": True, "color": "darkblue"},
            "image_placeholder": {"x": 8.0, "y": 1.5, "w": 4.833, "h": 4.5, "type": "picture"},
            "logo": {"x": 11.5, "y": 0.3, "w": 1.5, "h": 0.6},
            "page_info": {"x": 0.3, "y": 0.3, "w": 4.0, "h": 0.4, "font_size": 10, "color": "darkblue"},
        }
    },
}


def create_master_template(output_path: Path):
    """
    Erstellt ein FehrAdvice PPTX Master-Template

    Da python-pptx keine direkte Manipulation von Slide Masters erlaubt,
    erstellen wir ein Template mit vordefinierten Beispiel-Slides für jedes Layout.
    Diese können dann als Vorlage in PowerPoint verwendet werden.
    """
    print("Creating FehrAdvice Master Template...")

    # Neue Präsentation erstellen
    prs = Presentation()

    # Slide-Dimensionen setzen (16:9)
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Für jedes Layout einen Beispiel-Slide erstellen
    for layout_id, layout_def in LAYOUTS.items():
        print(f"  Creating layout: {layout_def['name']}")
        create_example_slide(prs, layout_id, layout_def)

    # Speichern
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"\n✅ Master Template saved to: {output_path}")

    return prs


def create_example_slide(prs: Presentation, layout_id: str, layout_def: dict):
    """Erstellt einen Beispiel-Slide für ein Layout"""

    # Blank layout verwenden
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    elements = layout_def["elements"]

    # Hintergrund
    if "background" in elements:
        bg = elements["background"]
        if bg.get("type") == "solid" and bg.get("color"):
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = COLORS[bg["color"]]

    # Header-Linie
    if "header_line" in elements:
        el = elements["header_line"]
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(el["x"]), Inches(el["y"]),
            Inches(el["w"]), Inches(el["h"])
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS[el.get("color", "darkblue")]
        shape.line.fill.background()

    # Titel
    if "title" in elements:
        el = elements["title"]
        add_text_box(slide, el, f"[{layout_def['name']}]", is_title=True)

    # Untertitel
    if "subtitle" in elements:
        el = elements["subtitle"]
        add_text_box(slide, el, "[Untertitel hier]")

    # Content
    if "content" in elements:
        el = elements["content"]
        add_text_box(slide, el, "[Inhalt hier einfügen]\n\n• Punkt 1\n• Punkt 2\n• Punkt 3")

    # Left/Right Content für Two-Column
    if "left_content" in elements:
        el = elements["left_content"]
        add_text_box(slide, el, "[Linke Spalte]\n\nText hier einfügen...")

    if "right_content" in elements:
        el = elements["right_content"]
        add_placeholder_box(slide, el, "Bild/Chart")

    # Chart Placeholder
    if "chart_placeholder" in elements:
        el = elements["chart_placeholder"]
        add_placeholder_box(slide, el, "Chart hier einfügen")

    # KPI Boxes
    for key in ["kpi_1", "kpi_2", "kpi_3", "kpi_4"]:
        if key in elements:
            el = elements[key]
            add_kpi_box(slide, el, f"KPI {key[-1]}")

    # Step Boxes für Process
    for key in ["step_1", "step_2", "step_3", "step_4", "step_5"]:
        if key in elements:
            el = elements[key]
            add_step_box(slide, el)

    # Timeline Line
    if "timeline_line" in elements:
        el = elements["timeline_line"]
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(el["x"]), Inches(el["y"]),
            Inches(el["w"]), Inches(el["h"])
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS[el.get("color", "lightblue")]
        shape.line.fill.background()

    # Quote
    if "quote" in elements:
        el = elements["quote"]
        add_text_box(slide, el, "«Zitat oder wichtige\nAussage hier.»")

    # Image Placeholder
    if "image_placeholder" in elements:
        el = elements["image_placeholder"]
        add_placeholder_box(slide, el, "Bild")

    # Logo (oben rechts)
    if "logo" in elements:
        el = elements["logo"]
        # Check if background is dark (for white logo)
        bg = elements.get("background", {})
        dark_bg = bg.get("color") == "darkblue" or bg.get("type") == "image"
        add_logo_placeholder(slide, el, dark_background=dark_bg)

    # Footer
    if "footer" in elements:
        el = elements["footer"]
        add_text_box(slide, el, "FehrAdvice & Partners AG | Präsentationstitel | Autor | Datum | Seite X")

    # Page Info (für Section Slides)
    if "page_info" in elements:
        el = elements["page_info"]
        add_text_box(slide, el, "X | FehrAdvice & Partners AG | CI Guide | 07.10.23")


def add_text_box(slide, el: dict, text: str, is_title: bool = False):
    """Fügt eine Textbox hinzu"""
    textbox = slide.shapes.add_textbox(
        Inches(el["x"]), Inches(el["y"]),
        Inches(el["w"]), Inches(el["h"])
    )
    tf = textbox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(el.get("font_size", 16))
    p.font.bold = el.get("bold", False)
    p.font.color.rgb = COLORS[el.get("color", "darkgray")]

    if el.get("align") == "center":
        p.alignment = PP_ALIGN.CENTER


def add_placeholder_box(slide, el: dict, label: str):
    """Fügt eine Platzhalter-Box hinzu"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(el["x"]), Inches(el["y"]),
        Inches(el["w"]), Inches(el["h"])
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["lightgray"]
    shape.line.color.rgb = COLORS["lightblue"]
    shape.line.width = Pt(2)

    # Label in der Mitte
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"[{label}]"
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS["darkgray"]
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(el["h"] * 72 / 2 - 10)


def add_kpi_box(slide, el: dict, label: str):
    """Fügt eine KPI-Box hinzu"""
    # Box
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(el["x"]), Inches(el["y"]),
        Inches(el["w"]), Inches(el["h"])
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["darkblue"]
    shape.line.fill.background()

    # KPI Value
    value_box = slide.shapes.add_textbox(
        Inches(el["x"]), Inches(el["y"] + 0.5),
        Inches(el["w"]), Inches(1.0)
    )
    tf = value_box.text_frame
    p = tf.paragraphs[0]
    p.text = "4.531 | 62%"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS["white"]
    p.alignment = PP_ALIGN.CENTER

    # KPI Label
    label_box = slide.shapes.add_textbox(
        Inches(el["x"]), Inches(el["y"] + 1.5),
        Inches(el["w"]), Inches(0.8)
    )
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS["white"]
    p.alignment = PP_ALIGN.CENTER


def add_step_box(slide, el: dict):
    """Fügt eine Step-Box für Process-Layout hinzu"""
    # Nummer-Kreis
    number = el.get("number", "01")
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(el["x"] + el["w"]/2 - 0.3), Inches(el["y"] + 2.5),
        Inches(0.6), Inches(0.6)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLORS["darkblue"]
    circle.line.fill.background()

    # Nummer
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS["white"]
    p.alignment = PP_ALIGN.CENTER

    # Title Box
    title_box = slide.shapes.add_textbox(
        Inches(el["x"]), Inches(el["y"]),
        Inches(el["w"]), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Title {number}"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS["darkblue"]
    p.alignment = PP_ALIGN.CENTER


def add_logo_placeholder(slide, el: dict, dark_background: bool = False):
    """Fügt FehrAdvice Logo hinzu

    Args:
        slide: PowerPoint slide
        el: Element definition with x, y, w, h
        dark_background: True = white logo, False = color logo
    """
    logo_path = LOGO_WHITE if dark_background else LOGO_COLOR

    if logo_path.exists():
        slide.shapes.add_picture(
            str(logo_path),
            Inches(el["x"]), Inches(el["y"]),
            height=Inches(el["h"])
        )
    else:
        # Fallback: placeholder box
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(el["x"]), Inches(el["y"]),
            Inches(el["w"]), Inches(el["h"])
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS["lightgray"]
        shape.line.color.rgb = COLORS["darkblue"]
        shape.line.width = Pt(1)

        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = "LOGO"
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS["darkblue"]
        p.alignment = PP_ALIGN.CENTER


def main():
    """Main entry point"""
    import argparse

    parser = argparse.ArgumentParser(
        description="Create FehrAdvice PPTX Master Template"
    )
    parser.add_argument(
        "--output", "-o",
        type=str,
        default="templates/pptx/FehrAdvice-Master.pptx",
        help="Output path for master template"
    )
    parser.add_argument(
        "--list-layouts",
        action="store_true",
        help="List available layouts"
    )

    args = parser.parse_args()

    if args.list_layouts:
        print("\nAvailable Layouts:")
        print("=" * 50)
        for layout_id, layout_def in LAYOUTS.items():
            print(f"\n{layout_id}:")
            print(f"  Name: {layout_def['name']}")
            print(f"  Description: {layout_def['description']}")
        return

    # Template erstellen
    output_path = Path(args.output)
    if not output_path.is_absolute():
        # Relativ zum Repo-Root
        repo_root = Path(__file__).parent.parent
        output_path = repo_root / output_path

    create_master_template(output_path)

    print("\n" + "=" * 60)
    print("NEXT STEPS:")
    print("=" * 60)
    print("""
1. Open the template in PowerPoint
2. View > Slide Master
3. Customize the layouts as needed:
   - Add FehrAdvice logo
   - Adjust fonts (Roboto, Open Sans)
   - Fine-tune colors and positions
4. Save as .potx (PowerPoint Template)

The generate_pptx.py script will use these layouts.
""")


if __name__ == "__main__":
    main()
