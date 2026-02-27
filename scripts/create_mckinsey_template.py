#!/usr/bin/env python3
"""
McKinsey-Style Slide Master Template Generator
===============================================

Erstellt ein PPTX Master-Template im McKinsey-Stil mit FehrAdvice-Farben.

McKinsey Design-Prinzipien (aus Slideworks & SlideModel):
1. Action Titles - Titel ist die Kernaussage als Satz
2. Pyramid Principle - Fazit zuerst, dann Begründung
3. Minimalistisch - Viel Weissraum, keine Animationen
4. Kontrast - Wechsel zwischen weiss und dunkelblau
5. Konsistente Schriften - Serif für Titel, Sans-Serif für Body
6. Quellen & Seitenzahl - Immer unten

FehrAdvice-Anpassungen:
- Dunkelblau: #024079 (statt McKinsey Blue)
- Hellblau: #549EDE (für Akzente)
- Roboto für Headlines (FehrAdvice Standard)
- Open Sans für Body Text

Version: 1.0
Date: January 2026

Sources:
- https://slideworks.io/resources/how-mckinsey-consultants-make-presentations
- https://slidemodel.com/mckinsey-presentation-structure/
"""

import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Error: python-pptx nicht installiert. Installiere mit: pip install python-pptx")
    sys.exit(1)

# =============================================================================
# FEHRADVICE COLORS (CI Guide)
# =============================================================================

COLORS = {
    "darkblue": RGBColor(0x02, 0x40, 0x79),      # #024079 - Primary
    "lightblue": RGBColor(0x54, 0x9E, 0xDE),     # #549EDE - Accent
    "darkgray": RGBColor(0x25, 0x21, 0x2A),      # #25212A - Text
    "lightgray": RGBColor(0xF3, 0xF5, 0xF7),     # #F3F5F7 - Background
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "mediumgray": RGBColor(0x66, 0x66, 0x66),    # For secondary text
}

# =============================================================================
# MCKINSEY-STYLE DIMENSIONS
# =============================================================================

SLIDE_WIDTH = Inches(13.333)   # 16:9
SLIDE_HEIGHT = Inches(7.5)

# Margins (McKinsey: Never go outside margins)
MARGIN_LEFT = Inches(0.5)
MARGIN_RIGHT = Inches(0.5)
MARGIN_TOP = Inches(0.4)
MARGIN_BOTTOM = Inches(0.5)

CONTENT_WIDTH = Inches(12.333)  # SLIDE_WIDTH - margins

# Logo paths
ASSETS_DIR = Path(__file__).parent.parent / "assets"
LOGO_WHITE = ASSETS_DIR / "logo" / "monochrome" / "fehradvice-logo-white.png"
LOGO_COLOR = ASSETS_DIR / "logo" / "color" / "fehradvice-logo-color-tagline.png"

# =============================================================================
# MCKINSEY LAYOUT DEFINITIONS
# =============================================================================

def create_mckinsey_template(output_path: Path):
    """Erstellt McKinsey-Style Template mit FehrAdvice-Farben"""

    print("Creating McKinsey-Style Template...")
    print("  Using FehrAdvice colors: #024079, #549EDE, #25212A")

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Slide 1: Title Slide (Dark Blue Background)
    create_title_slide(prs)

    # Slide 2: Agenda/Contents Slide
    create_agenda_slide(prs)

    # Slide 3: Section Divider (Dark Blue)
    create_section_divider(prs)

    # Slide 4: Standard Content Slide (White)
    create_content_slide(prs)

    # Slide 5: Two-Column Slide
    create_two_column_slide(prs)

    # Slide 6: Chart Slide
    create_chart_slide(prs)

    # Slide 7: Key Takeaway Slide
    create_takeaway_slide(prs)

    # Slide 8: Recommendation Slide
    create_recommendation_slide(prs)

    # Slide 9: Sources/End Slide
    create_sources_slide(prs)

    # Save
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    print(f"\n✅ McKinsey-Style Template saved to: {output_path}")
    return prs


def add_logo(slide, x=Inches(11.3), y=Inches(0.15), dark_background=False):
    """Fügt FehrAdvice Logo hinzu (oben rechts)

    Args:
        slide: PowerPoint slide
        x, y: Position
        dark_background: True = weisses Logo, False = farbiges Logo
    """
    logo_path = LOGO_WHITE if dark_background else LOGO_COLOR

    if logo_path.exists():
        # Logo einfügen (Höhe fixiert, Breite proportional)
        slide.shapes.add_picture(
            str(logo_path),
            x, y,
            height=Inches(0.5)
        )
    else:
        # Fallback: Text-Platzhalter
        color = "white" if dark_background else "darkblue"
        logo = slide.shapes.add_textbox(x, y, Inches(1.5), Inches(0.4))
        tf = logo.text_frame
        p = tf.paragraphs[0]
        p.text = "FEHR\nADVICE"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLORS[color]


def add_logo_placeholder(slide, x=Inches(11.5), y=Inches(0.2), color="darkblue"):
    """Legacy function - redirects to add_logo"""
    dark_background = (color == "white")
    add_logo(slide, x, y, dark_background)


def add_footer(slide, page_num: int, source: str = "", color="mediumgray"):
    """Fügt McKinsey-Style Footer hinzu (Quelle links, Seitenzahl rechts)"""
    # Source (left)
    if source:
        src_box = slide.shapes.add_textbox(
            MARGIN_LEFT, Inches(7.1),
            Inches(10), Inches(0.3)
        )
        tf = src_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Source: {source}"
        p.font.size = Pt(8)
        p.font.color.rgb = COLORS[color]

    # Page number (right)
    page_box = slide.shapes.add_textbox(
        Inches(12.5), Inches(7.1),
        Inches(0.5), Inches(0.3)
    )
    tf = page_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(page_num)
    p.font.size = Pt(9)
    p.font.color.rgb = COLORS[color]
    p.alignment = PP_ALIGN.RIGHT


def add_action_title(slide, title: str, subtitle: str = None, dark_bg: bool = False):
    """
    McKinsey Action Title - Die Kernaussage als Satz

    McKinsey-Regel: "The title tells you the conclusion of the slide"
    """
    title_color = "white" if dark_bg else "darkblue"

    # Main title (Action Title = Sentence)
    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT, MARGIN_TOP,
        CONTENT_WIDTH, Inches(0.8)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS[title_color]

    # Subtitle (optional elaboration)
    if subtitle:
        sub_box = slide.shapes.add_textbox(
            MARGIN_LEFT, Inches(1.1),
            CONTENT_WIDTH, Inches(0.5)
        )
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS["lightblue" if dark_bg else "mediumgray"]


def create_title_slide(prs: Presentation):
    """Slide 1: Title Slide (Dark Blue Background)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Dark blue background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS["darkblue"]
    bg.line.fill.background()

    # Logo (white on dark)
    add_logo_placeholder(slide, color="white")

    # Main Title
    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(2.5),
        CONTENT_WIDTH, Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "[Präsentationstitel als Action Statement]"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS["white"]

    # Subtitle
    sub_box = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(4.2),
        CONTENT_WIDTH, Inches(0.8)
    )
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "[Optionaler Untertitel mit Kontext]"
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS["lightblue"]

    # Meta info (date, author)
    meta_box = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(6.5),
        CONTENT_WIDTH, Inches(0.5)
    )
    tf = meta_box.text_frame
    p = tf.paragraphs[0]
    p.text = "FehrAdvice & Partners AG | Januar 2026"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS["lightblue"]


def create_agenda_slide(prs: Presentation):
    """Slide 2: Agenda/Contents"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_logo_placeholder(slide)
    add_action_title(slide, "Diese Präsentation beantwortet drei zentrale Fragen")

    # Divider line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MARGIN_LEFT, Inches(1.5),
        CONTENT_WIDTH, Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["lightblue"]
    line.line.fill.background()

    # Agenda items (numbered)
    items = [
        ("01", "Erste zentrale Frage oder Themenbereich"),
        ("02", "Zweite zentrale Frage oder Themenbereich"),
        ("03", "Dritte zentrale Frage oder Themenbereich"),
    ]

    for i, (num, text) in enumerate(items):
        y = Inches(2.2 + i * 1.5)

        # Number box
        num_box = slide.shapes.add_textbox(MARGIN_LEFT, y, Inches(0.8), Inches(1))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = COLORS["darkblue"]

        # Text
        txt_box = slide.shapes.add_textbox(Inches(1.5), y, Inches(10), Inches(1))
        tf = txt_box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(20)
        p.font.color.rgb = COLORS["darkgray"]

    add_footer(slide, 2)


def create_section_divider(prs: Presentation):
    """Slide 3: Section Divider (Dark Blue)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Dark blue background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS["darkblue"]
    bg.line.fill.background()

    add_logo_placeholder(slide, color="white")

    # Section number
    num_box = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(2.5),
        Inches(1), Inches(1)
    )
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = "01"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS["lightblue"]

    # Section title
    title_box = slide.shapes.add_textbox(
        Inches(1.8), Inches(2.6),
        Inches(10), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "[Abschnittstitel als Statement]"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS["white"]

    add_footer(slide, 3, color="lightblue")


def create_content_slide(prs: Presentation):
    """Slide 4: Standard Content (White Background)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_logo_placeholder(slide)
    add_action_title(
        slide,
        "Action Title: Die Hauptaussage dieser Folie als vollständiger Satz",
        "Optionaler Untertitel mit zusätzlichem Kontext oder Einordnung"
    )

    # Divider line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MARGIN_LEFT, Inches(1.7),
        CONTENT_WIDTH, Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["lightgray"]
    line.line.fill.background()

    # Content area
    content_box = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(2.0),
        CONTENT_WIDTH, Inches(4.5)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    bullets = [
        "Erster Hauptpunkt mit unterstützender Evidenz",
        "Zweiter Hauptpunkt mit Daten oder Beispiel",
        "Dritter Hauptpunkt mit Implikation",
    ]

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS["darkgray"]
        p.space_after = Pt(16)

    add_footer(slide, 4, "FehrAdvice Analyse")


def create_two_column_slide(prs: Presentation):
    """Slide 5: Two-Column Layout"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_logo_placeholder(slide)
    add_action_title(slide, "Zwei-Spalten-Layout für Vergleiche oder Text + Visual")

    # Divider line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MARGIN_LEFT, Inches(1.5),
        CONTENT_WIDTH, Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["lightgray"]
    line.line.fill.background()

    # Left column
    left_box = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(1.8),
        Inches(5.5), Inches(4.8)
    )
    tf = left_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Linke Spalte"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS["darkblue"]

    p2 = tf.add_paragraph()
    p2.text = "\n• Punkt 1\n• Punkt 2\n• Punkt 3"
    p2.font.size = Pt(16)
    p2.font.color.rgb = COLORS["darkgray"]

    # Right column (placeholder for visual)
    right_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(7), Inches(1.8),
        Inches(5.833), Inches(4.8)
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = COLORS["lightgray"]
    right_box.line.color.rgb = COLORS["lightblue"]
    right_box.line.width = Pt(1)

    # Placeholder text
    txt = slide.shapes.add_textbox(
        Inches(7), Inches(3.8),
        Inches(5.833), Inches(1)
    )
    tf = txt.text_frame
    p = tf.paragraphs[0]
    p.text = "[Chart / Visual / Bild]"
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS["mediumgray"]
    p.alignment = PP_ALIGN.CENTER

    add_footer(slide, 5, "FehrAdvice Analyse")


def create_chart_slide(prs: Presentation):
    """Slide 6: Chart Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_logo_placeholder(slide)
    add_action_title(slide, "Chart-Titel beschreibt die Kernaussage des Charts")

    # Chart placeholder
    chart_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MARGIN_LEFT, Inches(1.5),
        CONTENT_WIDTH, Inches(5.2)
    )
    chart_box.fill.solid()
    chart_box.fill.fore_color.rgb = COLORS["lightgray"]
    chart_box.line.color.rgb = COLORS["lightblue"]
    chart_box.line.width = Pt(1)

    # Placeholder text
    txt = slide.shapes.add_textbox(
        Inches(4), Inches(3.5),
        Inches(5), Inches(1)
    )
    tf = txt.text_frame
    p = tf.paragraphs[0]
    p.text = "[Chart hier einfügen]"
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS["mediumgray"]
    p.alignment = PP_ALIGN.CENTER

    add_footer(slide, 6, "Datenquelle hier angeben")


def create_takeaway_slide(prs: Presentation):
    """Slide 7: Key Takeaway Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_logo_placeholder(slide)
    add_action_title(slide, "Key Takeaways: Die drei wichtigsten Erkenntnisse")

    # Divider line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MARGIN_LEFT, Inches(1.5),
        CONTENT_WIDTH, Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["lightblue"]
    line.line.fill.background()

    # Three takeaway boxes
    takeaways = [
        ("1", "Erste Erkenntnis", "Unterstützende Details oder Evidenz"),
        ("2", "Zweite Erkenntnis", "Unterstützende Details oder Evidenz"),
        ("3", "Dritte Erkenntnis", "Unterstützende Details oder Evidenz"),
    ]

    box_width = Inches(3.8)
    start_x = MARGIN_LEFT

    for i, (num, title, detail) in enumerate(takeaways):
        x = Inches(0.5 + i * 4.2)
        y = Inches(2.0)

        # Box background
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y, box_width, Inches(4.2)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = COLORS["lightgray"]
        box.line.fill.background()

        # Number
        num_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), Inches(0.5), Inches(0.6))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = COLORS["darkblue"]

        # Title
        title_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.9), box_width - Inches(0.4), Inches(1))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS["darkblue"]

        # Detail
        detail_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1.8), box_width - Inches(0.4), Inches(2))
        tf = detail_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = detail
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS["darkgray"]

    add_footer(slide, 7)


def create_recommendation_slide(prs: Presentation):
    """Slide 8: Recommendation Slide (McKinsey-Style)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_logo_placeholder(slide)
    add_action_title(slide, "Empfehlungen: Konkrete nächste Schritte")

    # Divider line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MARGIN_LEFT, Inches(1.5),
        CONTENT_WIDTH, Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["lightblue"]
    line.line.fill.background()

    # Recommendations (McKinsey: use action verbs)
    recs = [
        ("Analysieren", "Durchführen einer vertieften Analyse der Zielgruppe"),
        ("Implementieren", "Einführung der empfohlenen Intervention in Pilotphase"),
        ("Messen", "Etablieren von KPIs zur Erfolgsmessung"),
        ("Skalieren", "Ausrollen auf weitere Bereiche nach erfolgreichem Pilot"),
    ]

    for i, (action, detail) in enumerate(recs):
        y = Inches(2.0 + i * 1.2)

        # Action word (bold)
        action_box = slide.shapes.add_textbox(MARGIN_LEFT, y, Inches(2), Inches(0.5))
        tf = action_box.text_frame
        p = tf.paragraphs[0]
        p.text = action
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS["darkblue"]

        # Detail
        detail_box = slide.shapes.add_textbox(Inches(2.8), y, Inches(10), Inches(0.8))
        tf = detail_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = detail
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS["darkgray"]

    add_footer(slide, 8)


def create_sources_slide(prs: Presentation):
    """Slide 9: Sources/End Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Dark blue background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS["darkblue"]
    bg.line.fill.background()

    add_logo_placeholder(slide, color="white")

    # "Quellen" title
    title_box = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(1.5),
        CONTENT_WIDTH, Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Quellen & Referenzen"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS["white"]

    # Sources list
    sources_box = slide.shapes.add_textbox(
        MARGIN_LEFT, Inches(2.5),
        Inches(7), Inches(4)
    )
    tf = sources_box.text_frame
    tf.word_wrap = True

    sources = [
        "• FehrAdvice Kontext-Datenbank (BCM2)",
        "• Model Registry & Theory Catalog",
        "• Wissenschaftliche Literatur (bcm_master.bib)",
        "• Session-Analyse",
    ]

    for i, source in enumerate(sources):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = source
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS["lightblue"]
        p.space_after = Pt(12)

    # Contact box
    contact_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.5), Inches(2.5),
        Inches(4.333), Inches(3)
    )
    contact_box.fill.solid()
    contact_box.fill.fore_color.rgb = COLORS["lightblue"]
    contact_box.line.fill.background()

    # Contact info
    contact_txt = slide.shapes.add_textbox(
        Inches(8.8), Inches(2.8),
        Inches(3.733), Inches(2.4)
    )
    tf = contact_txt.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "FehrAdvice & Partners AG"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS["darkblue"]

    p2 = tf.add_paragraph()
    p2.text = "\ninfo@fehradvice.com\nwww.fehradvice.com"
    p2.font.size = Pt(12)
    p2.font.color.rgb = COLORS["darkblue"]

    add_footer(slide, 9, color="lightblue")


def main():
    """Main entry point"""
    import argparse

    parser = argparse.ArgumentParser(
        description="Create McKinsey-Style PPTX Template with FehrAdvice Colors"
    )
    parser.add_argument(
        "--output", "-o",
        type=str,
        default="templates/pptx/FehrAdvice-McKinsey.pptx",
        help="Output path for template"
    )

    args = parser.parse_args()

    output_path = Path(args.output)
    if not output_path.is_absolute():
        repo_root = Path(__file__).parent.parent
        output_path = repo_root / output_path

    create_mckinsey_template(output_path)

    print("\n" + "=" * 60)
    print("MCKINSEY-STYLE DESIGN PRINCIPLES:")
    print("=" * 60)
    print("""
1. ACTION TITLES
   → Titel = Kernaussage als vollständiger Satz
   → NICHT: "Ergebnisse" → SONDERN: "Die Analyse zeigt drei Handlungsfelder"

2. PYRAMID PRINCIPLE
   → Fazit zuerst, dann Begründung
   → Executive Summary am Anfang, nicht am Ende

3. MINIMALISTISCH
   → Viel Weissraum
   → Keine Animationen, keine unnötigen Grafiken
   → Konsistente Schriftgrössen

4. KONTRAST
   → Wechsel zwischen weissen und dunkelblauen Folien
   → Strukturiert die Präsentation visuell

5. QUELLEN & SEITENZAHL
   → Immer unten auf jeder Folie
   → Schafft Glaubwürdigkeit

Sources:
- https://slideworks.io/resources/how-mckinsey-consultants-make-presentations
- https://slidemodel.com/mckinsey-presentation-structure/
""")


if __name__ == "__main__":
    main()
