#!/usr/bin/env python3
"""
Generate SOG Presentation - Original Design (1-to-1 from PDF)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

class OriginalCI:
    RED = RGBColor(0xC8, 0x10, 0x2E)
    BLACK = RGBColor(0x00, 0x00, 0x00)
    DARKGRAY = RGBColor(0x33, 0x33, 0x33)
    LIGHTGRAY = RGBColor(0xE8, 0xE8, 0xE8)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    FONT_HEADLINE = "Arial"
    FONT_BODY = "Arial"
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

def set_style(p, font, size, color, bold=False, italic=False):
    p.font.name = font
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic

def add_title(slide, text, top=Inches(0.3), size=36):
    box = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.333), Inches(1.0))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    set_style(p, OriginalCI.FONT_HEADLINE, size, OriginalCI.BLACK, bold=True)
    return box

def add_status_badge(slide):
    badge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.0), Inches(6.7), Inches(3.0), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = OriginalCI.RED
    badge.line.fill.background()
    tf = badge.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = "Status: Strategisches Briefing"
    set_style(p, OriginalCI.FONT_BODY, 12, OriginalCI.WHITE)

def slide_01(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(2))
    tf = title.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "Ordnung im digitalen Raum:"
    set_style(p1, OriginalCI.FONT_HEADLINE, 44, OriginalCI.BLACK, bold=True)
    p2 = tf.add_paragraph()
    p2.text = "Das Social Media Ordnungs-Gesetz (SOG)"
    set_style(p2, OriginalCI.FONT_HEADLINE, 44, OriginalCI.BLACK, bold=True)
    sub = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12), Inches(0.8))
    p = sub.text_frame.paragraphs[0]
    p.text = "Strategisches Briefing zur Kommunikationsoffensive."
    set_style(p, OriginalCI.FONT_BODY, 24, OriginalCI.DARKGRAY)
    desc = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12), Inches(2))
    p = desc.text_frame.paragraphs[0]
    p.text = "[ Visualisierung: Rot-schwarze Punkte von Chaos zu Ordnung ]"
    p.alignment = PP_ALIGN.CENTER
    set_style(p, OriginalCI.FONT_BODY, 14, OriginalCI.DARKGRAY, italic=True)
    add_status_badge(slide)
    return slide

def slide_02(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Management Summary: Wir schaffen\nFairness und Verantwortung.", size=36)
    cols = [
        ("SITUATION", "Der digitale Raum ist im Vergleich zu analogen Medien (TV, Radio) weitgehend unreguliert.", False),
        ("KOMPLIKATION", "Dies gefaehrdet Kinder durch mangelnden Schutz und die Demokratie durch Desinformation und Polarisierung.", False),
        ("LOESUNG", "Das SOG (Social Media Ordnungs-Gesetz). Kein Verbot, sondern die Einfuehrung von fairen Spielregeln, wie sie offline laengst gelten.", True)
    ]
    for i, (header, text, is_red) in enumerate(cols):
        x = Inches(0.5 + i * 4.2)
        hdr = slide.shapes.add_textbox(x, Inches(2.0), Inches(4), Inches(0.5))
        p = hdr.text_frame.paragraphs[0]
        p.text = header
        p.alignment = PP_ALIGN.CENTER
        set_style(p, OriginalCI.FONT_HEADLINE, 18, OriginalCI.BLACK, bold=True)
        txt = slide.shapes.add_textbox(x, Inches(2.8), Inches(3.8), Inches(2.5))
        tf = txt.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        set_style(p, OriginalCI.FONT_BODY, 16, OriginalCI.DARKGRAY)
    quote_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.8), Inches(13.333), Inches(1.2))
    quote_box.fill.solid()
    quote_box.fill.fore_color.rgb = OriginalCI.LIGHTGRAY
    quote_box.line.fill.background()
    qt = slide.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(12.333), Inches(0.8))
    p = qt.text_frame.paragraphs[0]
    p.text = "Diese Regierung schafft Ordnung. Wer Reichweite hat, traegt Verantwortung."
    p.alignment = PP_ALIGN.CENTER
    set_style(p, OriginalCI.FONT_BODY, 18, OriginalCI.BLACK, bold=True, italic=True)
    return slide

def slide_03(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1.2))
    tf = title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = 'Der strategische Rahmen: Von "Verbot" zu "Ordnung".'
    set_style(p, OriginalCI.FONT_HEADLINE, 32, OriginalCI.BLACK, italic=True)
    p2 = tf.add_paragraph()
    p2.text = "Wir verlassen die Defensive und argumentieren staatstragend."
    set_style(p2, OriginalCI.FONT_BODY, 24, OriginalCI.BLACK, bold=True)
    left_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.8), Inches(6.5), Inches(4.2))
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = OriginalCI.WHITE
    left_bg.line.fill.background()
    right_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(1.8), Inches(6.833), Inches(4.2))
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = OriginalCI.LIGHTGRAY
    right_bg.line.fill.background()
    left_hdr = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(5.5), Inches(0.6))
    p = left_hdr.text_frame.paragraphs[0]
    p.text = "NICHT MEHR SAGEN"
    set_style(p, OriginalCI.FONT_HEADLINE, 16, OriginalCI.DARKGRAY, bold=True)
    old_terms = ["Verbot / Einschraenkung", "Ich / Mein Ministerium", "Strafen", "Defensive Erklaerungen"]
    for i, term in enumerate(old_terms):
        box = slide.shapes.add_textbox(Inches(0.5), Inches(2.7 + i * 0.8), Inches(5.5), Inches(0.6))
        p = box.text_frame.paragraphs[0]
        p.text = "--- " + term + " ---"
        set_style(p, OriginalCI.FONT_BODY, 20, OriginalCI.DARKGRAY)
    right_hdr = slide.shapes.add_textbox(Inches(7), Inches(2.0), Inches(5.5), Inches(0.6))
    p = right_hdr.text_frame.paragraphs[0]
    p.text = "NEU: STRATEGISCHE TONALITAET"
    set_style(p, OriginalCI.FONT_HEADLINE, 16, OriginalCI.BLACK, bold=True)
    new_terms = ["Ordnung / Regeln / Schutz", "Diese Regierung / Wir als Regierung", "Verantwortung einfordern / Konsequenzen", "Gleiche Spielregeln fuer alle"]
    for i, term in enumerate(new_terms):
        box = slide.shapes.add_textbox(Inches(7), Inches(2.7 + i * 0.8), Inches(5.5), Inches(0.6))
        p = box.text_frame.paragraphs[0]
        p.text = "-> " + term
        set_style(p, OriginalCI.FONT_BODY, 20, OriginalCI.BLACK, bold=True)
    quote = slide.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(12.333), Inches(0.8))
    p = quote.text_frame.paragraphs[0]
    p.text = "Wir rechtfertigen uns nicht fuer Regeln. Wir schaffen sie."
    p.alignment = PP_ALIGN.CENTER
    set_style(p, OriginalCI.FONT_HEADLINE, 24, OriginalCI.RED, bold=True)
    return slide

def slide_04(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Vier Saeulen fuer die digitale Souveraenitaet.", size=36)
    pillars = [
        ("Altersgerechter\nZugang", "Schutz fuer Kinder und Jugendliche"),
        ("Plattform-\nHaftung", "Wer verdient, haftet"),
        ("Medien-\nkompetenz", "Demokratie braucht muendige Buerger:innen"),
        ("Transparenz", "Gleiche Spielregeln fuer alle Medien")
    ]
    for i, (title_text, desc) in enumerate(pillars):
        x = Inches(0.5 + i * 3.2)
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.8), Inches(2.9), Inches(1.5))
        header.fill.solid()
        header.fill.fore_color.rgb = OriginalCI.RED
        header.line.fill.background()
        body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(3.3), Inches(2.9), Inches(2.5))
        body.fill.solid()
        body.fill.fore_color.rgb = OriginalCI.LIGHTGRAY
        body.line.fill.background()
        ttl = slide.shapes.add_textbox(x + Inches(0.1), Inches(3.5), Inches(2.7), Inches(1))
        tf = ttl.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.alignment = PP_ALIGN.CENTER
        set_style(p, OriginalCI.FONT_HEADLINE, 20, OriginalCI.BLACK, bold=True)
        dsc = slide.shapes.add_textbox(x + Inches(0.1), Inches(4.6), Inches(2.7), Inches(1))
        tf = dsc.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.alignment = PP_ALIGN.CENTER
        set_style(p, OriginalCI.FONT_BODY, 14, OriginalCI.DARKGRAY)
    return slide

def slide_05(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Saeule 1: Konsequenter Schutz fuer unsere Kinder.", size=32)
    quote = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(4))
    tf = quote.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Unter 15 Jahren brauchen Kinder Schutz. Das ist konsequenter Jugendschutz - so wie bei Filmen, so wie bei Alkohol."
    set_style(p, OriginalCI.FONT_HEADLINE, 28, OriginalCI.BLACK, bold=True, italic=True)
    strat = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7), Inches(4.5), Inches(5.5), Inches(1.8))
    strat.fill.solid()
    strat.fill.fore_color.rgb = OriginalCI.LIGHTGRAY
    strat.line.fill.background()
    strat_txt = slide.shapes.add_textbox(Inches(7.2), Inches(4.6), Inches(5.1), Inches(1.6))
    tf = strat_txt.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Strategischer Standpunkt:"
    set_style(p, OriginalCI.FONT_HEADLINE, 16, OriginalCI.BLACK, bold=True)
    p2 = tf.add_paragraph()
    p2.text = "Wenn Europa nicht schnell genug ist, gehen wir national voran. Das ist Ordnung schaffen. Nicht zusehen."
    set_style(p2, OriginalCI.FONT_BODY, 14, OriginalCI.DARKGRAY)
    return slide

def slide_06(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Saeule 2: Verantwortung und Haftung.", size=36)
    sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.333), Inches(1))
    p = sub.text_frame.paragraphs[0]
    p.text = "Wer Reichweite hat, traegt Verantwortung.\nWer verdient, haftet."
    p.alignment = PP_ALIGN.CENTER
    set_style(p, OriginalCI.FONT_HEADLINE, 24, OriginalCI.BLACK, bold=True)
    scale = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    p = scale.text_frame.paragraphs[0]
    p.text = "[ Visualisierung: Rote Waage mit Gewinn vs Haftung (DSA Rules) ]"
    p.alignment = PP_ALIGN.CENTER
    set_style(p, OriginalCI.FONT_BODY, 16, OriginalCI.DARKGRAY, italic=True)
    left = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(5.5), Inches(1.5))
    tf = left.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Der Mechanismus"
    set_style(p, OriginalCI.FONT_HEADLINE, 20, OriginalCI.BLACK, bold=True)
    p2 = tf.add_paragraph()
    p2.text = "Der DSA (Digital Services Act) gibt den Rahmen. Strafen bis zu 6% des Jahresumsatzes."
    set_style(p2, OriginalCI.FONT_BODY, 16, OriginalCI.DARKGRAY)
    right = slide.shapes.add_textbox(Inches(7), Inches(4.5), Inches(5.5), Inches(1.5))
    tf = right.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Das Framing"
    set_style(p, OriginalCI.FONT_HEADLINE, 20, OriginalCI.BLACK, bold=True)
    p2 = tf.add_paragraph()
    p2.text = "Im Fernsehen und Radio gelten Regeln. Das muss auch fuer Social Media gelten. Gleiche Spielregeln fuer alle Medien sind keine Einschraenkung - das ist Fairness."
    set_style(p2, OriginalCI.FONT_BODY, 16, OriginalCI.DARKGRAY)
    return slide

def slide_07(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Saeule 3: Medienkompetenz ist Demokratiepolitik.", size=32)
    ctx = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(5.8), Inches(1.5))
    ctx.fill.solid()
    ctx.fill.fore_color.rgb = OriginalCI.LIGHTGRAY
    ctx.line.fill.background()
    ctx_txt = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(5.4), Inches(1.3))
    tf = ctx_txt.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Kontext: Junge Menschen informieren sich ueber Plattformen. Dort verschwimmen Journalismus und Desinformation."
    set_style(p, OriginalCI.FONT_BODY, 16, OriginalCI.DARKGRAY)
    goal = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.8), Inches(1.5))
    goal.fill.solid()
    goal.fill.fore_color.rgb = OriginalCI.LIGHTGRAY
    goal.line.fill.background()
    goal_txt = slide.shapes.add_textbox(Inches(7), Inches(1.6), Inches(5.4), Inches(1.3))
    tf = goal_txt.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Ziel: Demokratie braucht muendige Buerger:innen, die Quellen bewerten koennen."
    set_style(p, OriginalCI.FONT_HEADLINE, 18, OriginalCI.RED, bold=True)
    stmt = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(8), Inches(1))
    p = stmt.text_frame.paragraphs[0]
    p.text = "Diese Regierung wird nicht nur darueber reden.\nWir bauen die Struktur auf."
    set_style(p, OriginalCI.FONT_HEADLINE, 24, OriginalCI.BLACK, bold=True)
    timeline = slide.shapes.add_textbox(Inches(9), Inches(4.5), Inches(3.5), Inches(1.5))
    p = timeline.text_frame.paragraphs[0]
    p.text = "Start der Massnahmen:\n2026."
    p.alignment = PP_ALIGN.CENTER
    set_style(p, OriginalCI.FONT_HEADLINE, 24, OriginalCI.BLACK, bold=True)
    return slide

def slide_08(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Saeule 4: Transparenz und Algorithmen.", size=36)
    sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.333), Inches(0.6))
    p = sub.text_frame.paragraphs[0]
    p.text = "Keine Sonderrechte fuer Big Tech."
    set_style(p, OriginalCI.FONT_HEADLINE, 24, OriginalCI.BLACK, bold=True)
    left_hdr = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(5.5), Inches(0.8))
    p = left_hdr.text_frame.paragraphs[0]
    p.text = "Traditionelle Medien\n(TV/Print)"
    set_style(p, OriginalCI.FONT_HEADLINE, 20, OriginalCI.BLACK, bold=True)
    left_txt = slide.shapes.add_textbox(Inches(0.5), Inches(2.9), Inches(5.5), Inches(1.2))
    tf = left_txt.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Fake News verboten.\nBetruegerische Werbung verboten."
    set_style(p, OriginalCI.FONT_BODY, 18, OriginalCI.DARKGRAY)
    right_hdr = slide.shapes.add_textbox(Inches(7), Inches(2.0), Inches(5.5), Inches(0.8))
    p = right_hdr.text_frame.paragraphs[0]
    p.text = "Social Media\nPlattformen"
    set_style(p, OriginalCI.FONT_HEADLINE, 20, OriginalCI.BLACK, bold=True)
    right_txt = slide.shapes.add_textbox(Inches(7), Inches(2.9), Inches(5.5), Inches(1.2))
    tf = right_txt.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Verdienen Milliarden mit Polarisierung, Radikalisierung und Suchtverhalten."
    set_style(p, OriginalCI.FONT_BODY, 18, OriginalCI.DARKGRAY)
    quote_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(5.0), Inches(12.333), Inches(1.2))
    quote_box.fill.solid()
    quote_box.fill.fore_color.rgb = OriginalCI.WHITE
    quote_box.line.color.rgb = OriginalCI.RED
    quote_box.line.width = Pt(3)
    quote = slide.shapes.add_textbox(Inches(0.7), Inches(5.2), Inches(12), Inches(0.8))
    p = quote.text_frame.paragraphs[0]
    p.text = "All das ist klassischen Medien verboten. Warum sollten wir es Social Media erlauben? Gleiche Spielregeln fuer alle."
    p.alignment = PP_ALIGN.CENTER
    set_style(p, OriginalCI.FONT_HEADLINE, 18, OriginalCI.BLACK, bold=True)
    return slide

def slide_09(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    p = title.text_frame.paragraphs[0]
    p.text = 'Counter-Narrativ: Die "Zensur"-Luege der FPOE.'
    set_style(p, OriginalCI.FONT_HEADLINE, 32, OriginalCI.BLACK, italic=True)
    bubble = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(4.5), Inches(2.5))
    bubble.fill.solid()
    bubble.fill.fore_color.rgb = OriginalCI.WHITE
    bubble.line.color.rgb = OriginalCI.DARKGRAY
    bubble_txt = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(4), Inches(1.5))
    p = bubble_txt.text_frame.paragraphs[0]
    p.text = "Die FPOE ruft:\nZensur!"
    p.alignment = PP_ALIGN.CENTER
    set_style(p, OriginalCI.FONT_HEADLINE, 28, OriginalCI.BLACK, bold=True)
    ans_hdr = slide.shapes.add_textbox(Inches(6), Inches(1.3), Inches(6.5), Inches(0.6))
    p = ans_hdr.text_frame.paragraphs[0]
    p.text = "Die strategische Antwort"
    set_style(p, OriginalCI.FONT_HEADLINE, 22, OriginalCI.RED, bold=True)
    answers = [
        "Ich erinnere daran: Unter FPOE-Regierungsbeteiligung ist in Sachen Plattformregulierung nichts passiert. Null.",
        "Die rufen Zensur - und meinen: Keine Regeln fuer ihre Freunde in den Tech-Konzernen.",
        "Im Fernsehen gelten Regeln. Im Radio gelten Regeln. Warum soll ausgerechnet TikTok keine Regeln haben?"
    ]
    for i, ans in enumerate(answers):
        box = slide.shapes.add_textbox(Inches(6), Inches(2.0 + i * 1.1), Inches(6.5), Inches(1))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "* " + ans
        set_style(p, OriginalCI.FONT_BODY, 14, OriginalCI.BLACK)
    stmt = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12.333), Inches(0.8))
    p = stmt.text_frame.paragraphs[0]
    p.text = "Das ist keine Zensur. Das ist Ordnung."
    p.alignment = PP_ALIGN.CENTER
    set_style(p, OriginalCI.FONT_HEADLINE, 28, OriginalCI.RED, bold=True)
    return slide

def slide_10(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Antworten auf kritische Fragen der Umsetzung.", size=32)
    q1_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(12.333), Inches(2.2))
    q1_box.fill.solid()
    q1_box.fill.fore_color.rgb = OriginalCI.LIGHTGRAY
    q1_box.line.fill.background()
    q1_txt = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.8), Inches(2))
    tf = q1_txt.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Frage 1: Werden Kinder das Verbot nicht umgehen (VPN)?"
    set_style(p, OriginalCI.FONT_HEADLINE, 20, OriginalCI.RED, bold=True)
    p2 = tf.add_paragraph()
    p2.text = "Natuerlich werden manche es versuchen. Wie beim Alkohol. Aber die klare Grenze wirkt - sie schafft Bewusstsein. Wer weiss, dass etwas potenziell schaedlich ist, geht anders damit um."
    set_style(p2, OriginalCI.FONT_BODY, 16, OriginalCI.DARKGRAY, italic=True)
    q2_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.0), Inches(12.333), Inches(2.2))
    q2_box.fill.solid()
    q2_box.fill.fore_color.rgb = OriginalCI.LIGHTGRAY
    q2_box.line.fill.background()
    q2_txt = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(11.8), Inches(2))
    tf = q2_txt.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Frage 2: Sehen Experten ein Grundrechtsproblem (Datenschutz)?"
    set_style(p, OriginalCI.FONT_HEADLINE, 20, OriginalCI.RED, bold=True)
    p2 = tf.add_paragraph()
    p2.text = "Ein Eingriff ist zulaessig, wenn er verhaeltnismaessig ist. Alle Erkenntnisse zeigen: Der Schutzbedarf ist erheblich. Wir arbeiten daran, dass der Eingriff so gering wie moeglich ist - aber wirksam bleibt."
    set_style(p2, OriginalCI.FONT_BODY, 16, OriginalCI.DARKGRAY, italic=True)
    return slide

def slide_11(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1.5))
    p = title.text_frame.paragraphs[0]
    p.text = "Wir warten nicht.\nWir handeln."
    set_style(p, OriginalCI.FONT_HEADLINE, 40, OriginalCI.BLACK, bold=True)
    amb_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.5), Inches(5.5), Inches(1.5))
    amb_box.fill.solid()
    amb_box.fill.fore_color.rgb = OriginalCI.LIGHTGRAY
    amb_box.line.color.rgb = OriginalCI.RED
    amb_txt = slide.shapes.add_textbox(Inches(7.2), Inches(1.6), Inches(5.1), Inches(1.3))
    tf = amb_txt.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Strategische Ambition"
    set_style(p, OriginalCI.FONT_HEADLINE, 18, OriginalCI.RED, bold=True)
    p2 = tf.add_paragraph()
    p2.text = "Wenn Europa nicht schnell genug ist, gehen wir national voran."
    set_style(p2, OriginalCI.FONT_BODY, 14, OriginalCI.BLACK)
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(0.5), Inches(4.5), Inches(12.333), Inches(0.5))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = OriginalCI.RED
    arrow.line.fill.background()
    m1 = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(5), Inches(1.5))
    tf = m1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Bis Sommer"
    set_style(p, OriginalCI.FONT_HEADLINE, 20, OriginalCI.BLACK, bold=True)
    p2 = tf.add_paragraph()
    p2.text = "Vorlage des Gesetzesentwurfs (SOG). Diese Regierung schafft Ordnung im digitalen Raum."
    set_style(p2, OriginalCI.FONT_BODY, 14, OriginalCI.DARKGRAY)
    m2 = slide.shapes.add_textbox(Inches(8), Inches(5.2), Inches(4.5), Inches(1.5))
    tf = m2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "2026"
    set_style(p, OriginalCI.FONT_HEADLINE, 20, OriginalCI.BLACK, bold=True)
    p2 = tf.add_paragraph()
    p2.text = "Implementierung der Massnahmen zur Medienkompetenz."
    set_style(p2, OriginalCI.FONT_BODY, 14, OriginalCI.DARKGRAY)
    return slide

def slide_12(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1.2))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Kommunikation auf einen Blick:\nDie 5 Kernbotschaften."
    set_style(p, OriginalCI.FONT_HEADLINE, 32, OriginalCI.BLACK, bold=True)
    sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.333), Inches(0.5))
    p = sub.text_frame.paragraphs[0]
    p.text = "Cheat Sheet fuer Interviews und Pressetermine."
    set_style(p, OriginalCI.FONT_BODY, 18, OriginalCI.DARKGRAY)
    messages = [
        ("1.", "SOG Allgemein", "Diese Regierung schafft Ordnung im digitalen Raum."),
        ("2.", "Altersschutz", "Kinder schuetzen. So wie bei Filmen, so bei Social Media."),
        ("3.", "Plattformen", "Wer Reichweite hat, traegt Verantwortung."),
        ("4.", "Gegen FPOE", "Gleiche Spielregeln fuer alle Medien. Das ist keine Zensur - das ist Ordnung."),
        ("5.", "Geschwindigkeit", "Wenn Europa nicht schnell genug ist, gehen wir voran.")
    ]
    positions = [(Inches(0.5), Inches(2.2)), (Inches(4.7), Inches(2.2)), (Inches(8.9), Inches(2.2)), (Inches(0.5), Inches(4.5)), (Inches(4.7), Inches(4.5))]
    for i, (num, title_text, quote) in enumerate(messages):
        x, y = positions[i]
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.9), Inches(2.0))
        box.fill.solid()
        box.fill.fore_color.rgb = OriginalCI.LIGHTGRAY
        box.line.color.rgb = OriginalCI.DARKGRAY
        num_txt = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.1), Inches(0.8), Inches(0.6))
        p = num_txt.text_frame.paragraphs[0]
        p.text = num
        set_style(p, OriginalCI.FONT_HEADLINE, 28, OriginalCI.RED, bold=True)
        ttl = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.6), Inches(3.7), Inches(0.5))
        p = ttl.text_frame.paragraphs[0]
        p.text = title_text
        set_style(p, OriginalCI.FONT_HEADLINE, 16, OriginalCI.BLACK, bold=True)
        qt = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(1.1), Inches(3.7), Inches(0.8))
        tf = qt.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = quote
        set_style(p, OriginalCI.FONT_BODY, 12, OriginalCI.DARKGRAY)
    return slide

def main():
    print("=" * 60)
    print("Generating SOG Presentation - ORIGINAL Design")
    print("=" * 60)
    prs = Presentation()
    prs.slide_width = OriginalCI.SLIDE_WIDTH
    prs.slide_height = OriginalCI.SLIDE_HEIGHT
    print("\nGenerating slides...")
    slide_01(prs); print("  * Slide 1: Title")
    slide_02(prs); print("  * Slide 2: Management Summary")
    slide_03(prs); print("  * Slide 3: Strategic Frame")
    slide_04(prs); print("  * Slide 4: Four Pillars")
    slide_05(prs); print("  * Slide 5: Saeule 1")
    slide_06(prs); print("  * Slide 6: Saeule 2")
    slide_07(prs); print("  * Slide 7: Saeule 3")
    slide_08(prs); print("  * Slide 8: Saeule 4")
    slide_09(prs); print("  * Slide 9: Counter-Narrativ")
    slide_10(prs); print("  * Slide 10: FAQ")
    slide_11(prs); print("  * Slide 11: Timeline")
    slide_12(prs); print("  * Slide 12: Key Messages")
    output_dir = "outputs/presentations"
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "SOG_Original.pptx")
    prs.save(output_path)
    print("\n" + "=" * 60)
    print("PRESENTATION GENERATED: " + output_path)
    print("=" * 60)
    print("\nOriginal Design: Red #C8102E, Black, Gray")
    print("Slides: 12 (1-to-1 from PDF)")
    return output_path

if __name__ == "__main__":
    main()
