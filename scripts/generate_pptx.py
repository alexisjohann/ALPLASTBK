#!/usr/bin/env python3
"""
EBF PPTX Generator - Generiert PowerPoint aus Session-Daten
=============================================================

Verwendet:
- templates/pptx/fa-style.yaml (Corporate Style)
- templates/pptx/slide-types.yaml (Slide-Layouts)
- templates/pptx/8d-slide-mapping.yaml (Zielgruppen-Profile)
- templates/pptx/graphic-mapping.yaml (Chart-Definitionen)

Generiert native .pptx Dateien mit python-pptx.

Usage:
    python scripts/generate_pptx.py --session EBF-S-2026-01-26-COG-001 --audience board
    python scripts/generate_pptx.py --session EBF-S-2026-01-26-COG-001 --custom-8d 0.4,0.3,0.85,0.2,G2,0.9,0.4,0.6
"""

import argparse
import yaml
import re
from pathlib import Path
from datetime import datetime
from typing import Dict, Any, List, Optional, Tuple
from dataclasses import dataclass

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.dml import MSO_THEME_COLOR
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Warning: python-pptx nicht installiert. Installiere mit: pip install python-pptx")


# =============================================================================
# CONFIGURATION
# =============================================================================

BASE_DIR = Path(__file__).parent.parent
TEMPLATES_DIR = BASE_DIR / "templates" / "pptx"
ASSETS_DIR = BASE_DIR / "assets"
OUTPUTS_DIR = BASE_DIR / "outputs" / "sessions"
DATA_DIR = BASE_DIR / "data"

# Master Templates
MASTER_TEMPLATE = TEMPLATES_DIR / "FehrAdvice-Master.pptx"      # CI Guide Style
MCKINSEY_TEMPLATE = TEMPLATES_DIR / "FehrAdvice-McKinsey.pptx"  # McKinsey Style


@dataclass
class Profile8D:
    """8D-Profil für Zielgruppen-Anpassung"""
    D1_knowledge: float      # Wissen (0=Laie, 1=Expert)
    D2_proximity: float      # Nähe zum Fachgebiet
    D3_scope: float          # Reichweite (0=persönlich, 1=gesellschaftlich)
    D4_time: float           # Verfügbare Zeit
    D5_goal: str             # Kommunikationsziel (G1-G7)
    D6_context: float        # Kontext (0=intern, 1=extern)
    D7_emotion: float        # Emotion (0=sachlich, 1=emotional)
    D8_persistence: float    # Persistenz (0=kurzlebig, 1=archiv)


@dataclass
class SlideSpec:
    """Spezifikation für einen einzelnen Slide"""
    slide_type: str
    index: int
    title: Optional[str] = None
    content_sources: Optional[List[str]] = None
    graphics: Optional[List[str]] = None
    content: Optional[Dict[str, Any]] = None


# =============================================================================
# YAML LOADERS
# =============================================================================

def load_yaml(path: Path) -> Dict[str, Any]:
    """Lädt YAML-Datei"""
    with open(path, 'r', encoding='utf-8') as f:
        return yaml.safe_load(f)


def load_style() -> Dict[str, Any]:
    """Lädt FehrAdvice Style Guide"""
    return load_yaml(TEMPLATES_DIR / "fa-style.yaml")


def load_slide_types() -> Dict[str, Any]:
    """Lädt Slide-Type Definitionen"""
    return load_yaml(TEMPLATES_DIR / "slide-types.yaml")


def load_8d_mapping() -> Dict[str, Any]:
    """Lädt 8D → Slide Mapping"""
    return load_yaml(TEMPLATES_DIR / "8d-slide-mapping.yaml")


def load_graphic_mapping() -> Dict[str, Any]:
    """Lädt Graphic Mapping"""
    return load_yaml(TEMPLATES_DIR / "graphic-mapping.yaml")


def load_model_registry() -> Dict[str, Any]:
    """Lädt Model Registry"""
    return load_yaml(DATA_DIR / "model-registry.yaml")


def load_session_data(session_id: str) -> Dict[str, Any]:
    """Lädt Session-Daten aus verschiedenen Quellen"""
    session_dir = OUTPUTS_DIR / session_id

    # Versuche Session-Report zu laden
    session_data = {
        "session_id": session_id,
        "created": datetime.now().strftime("%d. %B %Y"),
    }

    # Model Registry durchsuchen
    model_registry = load_model_registry()
    models = model_registry.get("models", [])

    # Handle both list and dict formats
    if isinstance(models, list):
        for model_data in models:
            if model_data.get("created_in_session") == session_id:
                session_data["model"] = model_data
                session_data["model_id"] = model_data.get("id", "unknown")
                break
    elif isinstance(models, dict):
        for model_id, model_data in models.items():
            if model_data.get("session_ref") == session_id:
                session_data["model"] = model_data
                session_data["model_id"] = model_id
                break

    # Session-Report lesen (falls vorhanden)
    report_files = list(session_dir.glob("F1_*.md")) if session_dir.exists() else []
    if report_files:
        with open(report_files[0], 'r', encoding='utf-8') as f:
            content = f.read()
            session_data["report_content"] = content

            # Extrahiere Key Findings
            session_data["key_findings"] = extract_key_findings(content)
            session_data["implications"] = extract_implications(content)

    return session_data


def extract_key_findings(content: str) -> List[str]:
    """Extrahiert Key Findings aus Report"""
    findings = []

    # Suche nach Hauptergebnis-Sektion
    patterns = [
        r"##.*(?:Ergebnis|Result|Finding).*\n(.*?)(?=\n##|\Z)",
        r"\*\*Hauptergebnis.*?\*\*:?\s*(.*?)(?=\n\n|\Z)",
        r"Median.*?(\d+\.?\d*)\s*(?:Level|Stufe)",
    ]

    for pattern in patterns:
        matches = re.findall(pattern, content, re.IGNORECASE | re.DOTALL)
        for match in matches:
            if isinstance(match, str) and len(match) > 10:
                findings.append(match.strip()[:200])

    # Fallback
    if not findings:
        findings = [
            "Menschen denken im Median 1-2 Schritte voraus",
            "70% erreichen maximal Level 1 der strategischen Tiefe",
            "Komplexität ist der Haupttreiber für mangelndes Verständnis"
        ]

    return findings[:5]


def extract_implications(content: str) -> List[str]:
    """Extrahiert Implikationen aus Report"""
    implications = []

    # Suche nach Implikationen-Sektion
    pattern = r"##.*(?:Implikation|Implication).*\n(.*?)(?=\n##|\Z)"
    matches = re.findall(pattern, content, re.IGNORECASE | re.DOTALL)

    for match in matches:
        # Extrahiere Bullet Points
        bullets = re.findall(r"[-*]\s+(.*?)(?=\n[-*]|\n\n|\Z)", match)
        implications.extend([b.strip() for b in bullets if len(b) > 5])

    # Fallback
    if not implications:
        implications = [
            "Komplexität reduzieren: Weniger Elemente, klarere Strukturen",
            "Feedback verbessern: Konsequenzen sichtbar machen",
            "Verzögerungen minimieren: Zeitnahe Rückmeldungen",
            "Abstraktion vermeiden: Konkrete, greifbare Beispiele"
        ]

    return implications[:6]


# =============================================================================
# 8D PROFILE MANAGEMENT
# =============================================================================

def get_audience_profile(audience: str, mapping: Dict[str, Any]) -> Profile8D:
    """Holt 8D-Profil für eine Zielgruppe"""
    profiles = mapping.get("audience_profiles", {})

    if audience not in profiles:
        raise ValueError(f"Unbekannte Zielgruppe: {audience}. Verfügbar: {list(profiles.keys())}")

    p = profiles[audience]["profile_8d"]
    return Profile8D(
        D1_knowledge=p["D1_knowledge"],
        D2_proximity=p["D2_proximity"],
        D3_scope=p["D3_scope"],
        D4_time=p["D4_time"],
        D5_goal=p["D5_goal"],
        D6_context=p["D6_context"],
        D7_emotion=p["D7_emotion"],
        D8_persistence=p["D8_persistence"],
    )


def parse_custom_8d(custom_str: str) -> Profile8D:
    """Parst Custom 8D-String: '0.4,0.3,0.85,0.2,G2,0.9,0.4,0.6'"""
    parts = custom_str.split(",")
    if len(parts) != 8:
        raise ValueError(f"Custom 8D benötigt 8 Werte, erhalten: {len(parts)}")

    return Profile8D(
        D1_knowledge=float(parts[0]),
        D2_proximity=float(parts[1]),
        D3_scope=float(parts[2]),
        D4_time=float(parts[3]),
        D5_goal=parts[4].strip(),
        D6_context=float(parts[5]),
        D7_emotion=float(parts[6]),
        D8_persistence=float(parts[7]),
    )


# =============================================================================
# SLIDE ARCHITECTURE GENERATOR
# =============================================================================

def determine_slide_count(profile: Profile8D, mapping: Dict[str, Any]) -> int:
    """Bestimmt Slide-Anzahl aus D4 (Zeit)"""
    rules = mapping.get("slide_count_rules", {}).get("by_D4", [])

    for rule in rules:
        range_min, range_max = rule["range"]
        if range_min <= profile.D4_time < range_max:
            # Mittelwert des Bereichs
            slides = rule["slides"]
            return (slides[0] + slides[1]) // 2

    # Fallback
    return 10


def determine_slide_types(profile: Profile8D, mapping: Dict[str, Any]) -> List[str]:
    """Bestimmt Slide-Typen aus D5 (Ziel)"""
    goal_key = f"{profile.D5_goal}_{'informieren' if profile.D5_goal == 'G1' else 'handeln' if profile.D5_goal == 'G2' else 'ueberzeugen'}"

    # Mapping von D5 zu Slide-Flow
    goal_mapping = {
        "G1": "G1_informieren",
        "G2": "G2_handeln",
        "G3": "G3_ueberzeugen",
        "G4": "G4_unterhalten",
    }

    goal_config = mapping.get("slide_types_by_goal", {}).get(goal_mapping.get(profile.D5_goal, "G1_informieren"), {})

    flow = goal_config.get("slide_flow", ["TITLE", "EXECUTIVE_SUMMARY", "DATA_CHART", "IMPLICATIONS", "SOURCES"])

    # Parse flow (entferne ? und * Marker)
    parsed_flow = []
    for item in flow:
        slide_type = item.rstrip("?*")
        is_optional = item.endswith("?")
        is_multiple = item.endswith("*")

        parsed_flow.append({
            "type": slide_type,
            "optional": is_optional,
            "multiple": is_multiple
        })

    return parsed_flow


def determine_text_complexity(profile: Profile8D, mapping: Dict[str, Any]) -> Dict[str, Any]:
    """Bestimmt Text-Komplexität aus D1 (Wissen)"""
    rules = mapping.get("text_complexity_rules", {}).get("by_D1", [])

    for rule in rules:
        range_min, range_max = rule["range"]
        if range_min <= profile.D1_knowledge < range_max:
            return rule["settings"]

    # Fallback: moderate
    return {
        "font_sizes": {"title": 40, "body": 20, "caption": 14},
        "max_words_per_slide": 60,
        "bullet_depth": 2,
        "formulas": "simplified",
        "jargon": "explained",
    }


def determine_content_balance(profile: Profile8D, mapping: Dict[str, Any]) -> Dict[str, Any]:
    """Bestimmt Content-Balance aus D7 (Emotion)"""
    rules = mapping.get("content_balance_rules", {}).get("by_D7", [])

    for rule in rules:
        range_min, range_max = rule["range"]
        if range_min <= profile.D7_emotion < range_max:
            return {
                "style": rule["style"],
                "balance": rule["balance"],
                "preferences": rule["preferences"],
            }

    # Fallback: balanced
    return {
        "style": "balanced",
        "balance": {"charts_ratio": 0.45, "text_ratio": 0.4, "whitespace_ratio": 0.15},
        "preferences": {"chart_types": ["bar", "pie", "line"], "use_icons": True},
    }


def generate_slide_architecture(
    profile: Profile8D,
    session_data: Dict[str, Any],
    mapping: Dict[str, Any]
) -> List[SlideSpec]:
    """Generiert komplette Slide-Architektur aus 8D-Profil"""

    slide_count = determine_slide_count(profile, mapping)
    slide_flow = determine_slide_types(profile, mapping)
    text_complexity = determine_text_complexity(profile, mapping)
    content_balance = determine_content_balance(profile, mapping)

    slides = []
    current_index = 1

    model_data = session_data.get("model", {})

    for flow_item in slide_flow:
        slide_type = flow_item["type"]
        is_multiple = flow_item["multiple"]

        # Überspringe optionale Slides wenn Limit erreicht
        if current_index > slide_count and flow_item["optional"]:
            continue

        if slide_type == "TITLE":
            slides.append(SlideSpec(
                slide_type="TITLE",
                index=current_index,
                title=model_data.get("question_template", "Verstehen Menschen Interdependenzen?"),
                content={
                    "subtitle": f"EBF Session {session_data.get('session_id', '')}",
                    "author": "FehrAdvice & Partners AG",
                    "date": session_data.get("created", datetime.now().strftime("%d. %B %Y")),
                }
            ))

        elif slide_type == "EXECUTIVE_SUMMARY":
            findings = session_data.get("key_findings", [])
            slides.append(SlideSpec(
                slide_type="EXECUTIVE_SUMMARY",
                index=current_index,
                title="Executive Summary",
                content={
                    "key_finding": findings[0] if findings else "Hauptergebnis der Analyse",
                    "kpis": [
                        {"label": "Median Level", "value": "1.5"},
                        {"label": "L0+L1 Anteil", "value": "70%"},
                        {"label": "Konfidenz", "value": "Hoch"},
                    ],
                }
            ))

        elif slide_type == "DATA_CHART":
            # Generiere mehrere Chart-Slides wenn multiple
            chart_specs = [
                {"title": "Cognitive Hierarchy", "graphic": "cognitive_hierarchy_bar"},
                {"title": "Verständnis nach Komplexität", "graphic": "v_n_decay_line"},
                {"title": "Sensitivitätsanalyse", "graphic": "sensitivity_donut"},
            ]

            for i, spec in enumerate(chart_specs):
                if not is_multiple and i > 0:
                    break
                if current_index > slide_count:
                    break

                slides.append(SlideSpec(
                    slide_type="DATA_CHART",
                    index=current_index,
                    title=spec["title"],
                    graphics=[spec["graphic"]],
                    content={
                        "insight": f"Zentrale Erkenntnis zu {spec['title']}",
                    }
                ))
                current_index += 1
            continue  # Skip index increment at end

        elif slide_type == "IMPLICATIONS":
            implications = session_data.get("implications", [])
            slides.append(SlideSpec(
                slide_type="IMPLICATIONS",
                index=current_index,
                title="Was bedeutet das für Sie?",
                content={
                    "implications": implications[:4],
                    "icons": ["simplify", "visibility", "time", "person"],
                }
            ))

        elif slide_type == "CALL_TO_ACTION":
            slides.append(SlideSpec(
                slide_type="CALL_TO_ACTION",
                index=current_index,
                title="Nächste Schritte",
                content={
                    "main_message": "Komplexität reduzieren ist der wichtigste Hebel",
                    "action_items": [
                        "Systeme vereinfachen",
                        "Feedback-Schleifen einbauen",
                        "Abstraktion vermeiden",
                    ],
                }
            ))

        elif slide_type == "SOURCES":
            slides.append(SlideSpec(
                slide_type="SOURCES",
                index=current_index,
                title="Quellen & Kontakt",
                content={
                    "sources": [
                        "Camerer et al. (2004): Cognitive Hierarchy",
                        "Kahneman (2011): Thinking, Fast and Slow",
                        "Simon (1955): Bounded Rationality",
                    ],
                    "contact": {
                        "company": "FehrAdvice & Partners AG",
                        "email": "info@fehradvice.com",
                        "web": "www.fehradvice.com",
                    }
                }
            ))

        elif slide_type == "METHODOLOGY":
            slides.append(SlideSpec(
                slide_type="METHODOLOGY",
                index=current_index,
                title="Methodologie",
                content={
                    "model_name": model_data.get("name", "IDV-2.0"),
                    "description": "Formales Modell basierend auf Cognitive Hierarchy Theory",
                    "parameters": ["Basisverständnis κ", "Lernkoeffizient λ", "Komplexitätsexponent α"],
                }
            ))

        else:
            # Generischer Slide
            slides.append(SlideSpec(
                slide_type=slide_type,
                index=current_index,
                title=slide_type.replace("_", " ").title(),
            ))

        current_index += 1

        # Stoppe wenn Limit erreicht
        if current_index > slide_count:
            break

    return slides


# =============================================================================
# PPTX GENERATION
# =============================================================================

def hex_to_rgb(hex_color: str) -> RGBColor:
    """Konvertiert Hex-Farbe zu RGBColor"""
    hex_color = hex_color.lstrip('#')
    r, g, b = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    return RGBColor(r, g, b)


def add_title_slide(prs: Presentation, spec: SlideSpec, style: Dict[str, Any]):
    """Fügt Titel-Slide hinzu"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    colors = style["colors"]["primary"]
    typography = style["typography"]

    # Hintergrund
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = hex_to_rgb(colors["darkblue"]["hex"])
    background.line.fill.background()

    # Haupttitel
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5),
        Inches(12.333), Inches(2)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = spec.title or "Titel"
    title_para.font.size = Pt(typography["sizes"]["title_main"])
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER

    # Untertitel
    if spec.content and spec.content.get("subtitle"):
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.5),
            Inches(12.333), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = spec.content["subtitle"]
        subtitle_para.font.size = Pt(typography["sizes"]["title_subtitle"])
        subtitle_para.font.color.rgb = hex_to_rgb(colors["lightblue"]["hex"])
        subtitle_para.alignment = PP_ALIGN.CENTER

    # Datum & Autor
    if spec.content:
        meta_text = f"{spec.content.get('author', '')} | {spec.content.get('date', '')}"
        meta_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.5),
            Inches(12.333), Inches(0.5)
        )
        meta_frame = meta_box.text_frame
        meta_para = meta_frame.paragraphs[0]
        meta_para.text = meta_text
        meta_para.font.size = Pt(typography["sizes"]["title_meta"])
        meta_para.font.color.rgb = RGBColor(200, 200, 200)
        meta_para.alignment = PP_ALIGN.CENTER

    return slide


def add_executive_summary_slide(prs: Presentation, spec: SlideSpec, style: Dict[str, Any]):
    """Fügt Executive Summary Slide hinzu"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    colors = style["colors"]["primary"]
    typography = style["typography"]

    # Slide-Titel
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(12.333), Inches(1)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = spec.title or "Executive Summary"
    title_para.font.size = Pt(typography["sizes"]["slide_title"])
    title_para.font.bold = True
    title_para.font.color.rgb = hex_to_rgb(colors["darkblue"]["hex"])

    # Key Finding Box
    if spec.content and spec.content.get("key_finding"):
        finding_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1.5),
            Inches(12.333), Inches(1.5)
        )
        finding_box.fill.solid()
        finding_box.fill.fore_color.rgb = hex_to_rgb(colors["lightgray"]["hex"])
        finding_box.line.fill.background()

        # Text in Box
        finding_text = slide.shapes.add_textbox(
            Inches(0.7), Inches(1.7),
            Inches(11.933), Inches(1.1)
        )
        finding_frame = finding_text.text_frame
        finding_para = finding_frame.paragraphs[0]
        finding_para.text = spec.content["key_finding"]
        finding_para.font.size = Pt(typography["sizes"]["body_large"])
        finding_para.font.color.rgb = hex_to_rgb(colors["darkgray"]["hex"])

    # KPIs
    if spec.content and spec.content.get("kpis"):
        kpis = spec.content["kpis"]
        kpi_width = 3.5
        start_x = (13.333 - (len(kpis) * kpi_width)) / 2

        for i, kpi in enumerate(kpis):
            x = start_x + (i * kpi_width)

            # KPI Box
            kpi_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(3.5),
                Inches(3), Inches(2.5)
            )
            kpi_box.fill.solid()
            kpi_box.fill.fore_color.rgb = hex_to_rgb(colors["darkblue"]["hex"])
            kpi_box.line.fill.background()

            # KPI Value
            value_box = slide.shapes.add_textbox(
                Inches(x), Inches(3.8),
                Inches(3), Inches(1.2)
            )
            value_frame = value_box.text_frame
            value_para = value_frame.paragraphs[0]
            value_para.text = kpi["value"]
            value_para.font.size = Pt(36)
            value_para.font.bold = True
            value_para.font.color.rgb = RGBColor(255, 255, 255)
            value_para.alignment = PP_ALIGN.CENTER

            # KPI Label
            label_box = slide.shapes.add_textbox(
                Inches(x), Inches(5.0),
                Inches(3), Inches(0.8)
            )
            label_frame = label_box.text_frame
            label_para = label_frame.paragraphs[0]
            label_para.text = kpi["label"]
            label_para.font.size = Pt(14)
            label_para.font.color.rgb = hex_to_rgb(colors["lightblue"]["hex"])
            label_para.alignment = PP_ALIGN.CENTER

    return slide


def add_data_chart_slide(
    prs: Presentation,
    spec: SlideSpec,
    style: Dict[str, Any],
    session_id: str
):
    """Fügt Data Chart Slide hinzu"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    colors = style["colors"]["primary"]
    typography = style["typography"]

    # Slide-Titel
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(12.333), Inches(1)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = spec.title or "Datenanalyse"
    title_para.font.size = Pt(typography["sizes"]["slide_title"])
    title_para.font.bold = True
    title_para.font.color.rgb = hex_to_rgb(colors["darkblue"]["hex"])

    # Grafik einfügen (falls vorhanden)
    if spec.graphics:
        graphic_name = spec.graphics[0]
        graphic_path = ASSETS_DIR / "graphics" / "generated" / session_id / f"{graphic_name}.png"

        if graphic_path.exists():
            slide.shapes.add_picture(
                str(graphic_path),
                Inches(1), Inches(1.5),
                width=Inches(8),
                height=Inches(5)
            )
        else:
            # Placeholder-Box
            placeholder = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(1), Inches(1.5),
                Inches(8), Inches(5)
            )
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = hex_to_rgb(colors["lightgray"]["hex"])
            placeholder.line.color.rgb = hex_to_rgb(colors["darkblue"]["hex"])

            # Placeholder Text
            text_box = slide.shapes.add_textbox(
                Inches(1), Inches(3.5),
                Inches(8), Inches(1)
            )
            text_frame = text_box.text_frame
            text_para = text_frame.paragraphs[0]
            text_para.text = f"[Grafik: {graphic_name}]"
            text_para.font.size = Pt(18)
            text_para.font.color.rgb = hex_to_rgb(colors["darkgray"]["hex"])
            text_para.alignment = PP_ALIGN.CENTER

    # Insight-Text
    if spec.content and spec.content.get("insight"):
        insight_box = slide.shapes.add_textbox(
            Inches(9.5), Inches(2),
            Inches(3.333), Inches(4)
        )
        insight_frame = insight_box.text_frame
        insight_frame.word_wrap = True
        insight_para = insight_frame.paragraphs[0]
        insight_para.text = spec.content["insight"]
        insight_para.font.size = Pt(typography["sizes"]["body_normal"])
        insight_para.font.color.rgb = hex_to_rgb(colors["darkgray"]["hex"])

    return slide


def add_implications_slide(prs: Presentation, spec: SlideSpec, style: Dict[str, Any]):
    """Fügt Implications Slide hinzu"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    colors = style["colors"]["primary"]
    secondary = style["colors"]["secondary"]
    typography = style["typography"]

    # Slide-Titel
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(12.333), Inches(1)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = spec.title or "Implikationen"
    title_para.font.size = Pt(typography["sizes"]["slide_title"])
    title_para.font.bold = True
    title_para.font.color.rgb = hex_to_rgb(colors["darkblue"]["hex"])

    # Implikationen als Kacheln
    if spec.content and spec.content.get("implications"):
        implications = spec.content["implications"]

        # 2x2 Grid
        positions = [
            (0.5, 1.5), (6.667, 1.5),
            (0.5, 4.2), (6.667, 4.2)
        ]

        box_colors = [
            colors["darkblue"]["hex"],
            secondary["mint"]["hex"],
            secondary["lilac"]["hex"],
            secondary["ocher"]["hex"]
        ]

        for i, impl in enumerate(implications[:4]):
            x, y = positions[i]

            # Box
            impl_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y),
                Inches(6), Inches(2.5)
            )
            impl_box.fill.solid()
            impl_box.fill.fore_color.rgb = hex_to_rgb(box_colors[i])
            impl_box.line.fill.background()

            # Icon-Platzhalter (Nummer)
            icon_box = slide.shapes.add_textbox(
                Inches(x + 0.3), Inches(y + 0.3),
                Inches(1), Inches(1)
            )
            icon_frame = icon_box.text_frame
            icon_para = icon_frame.paragraphs[0]
            icon_para.text = str(i + 1)
            icon_para.font.size = Pt(36)
            icon_para.font.bold = True
            icon_para.font.color.rgb = RGBColor(255, 255, 255)

            # Text
            text_box = slide.shapes.add_textbox(
                Inches(x + 0.3), Inches(y + 1.2),
                Inches(5.4), Inches(1.2)
            )
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            text_para = text_frame.paragraphs[0]
            text_para.text = impl
            text_para.font.size = Pt(16)
            text_para.font.color.rgb = RGBColor(255, 255, 255) if i == 0 else hex_to_rgb(colors["darkgray"]["hex"])

    return slide


def add_call_to_action_slide(prs: Presentation, spec: SlideSpec, style: Dict[str, Any]):
    """Fügt Call-to-Action Slide hinzu"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    colors = style["colors"]["primary"]
    typography = style["typography"]

    # Slide-Titel
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(12.333), Inches(1)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = spec.title or "Nächste Schritte"
    title_para.font.size = Pt(typography["sizes"]["slide_title"])
    title_para.font.bold = True
    title_para.font.color.rgb = hex_to_rgb(colors["darkblue"]["hex"])

    # Main Message
    if spec.content and spec.content.get("main_message"):
        msg_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1.5),
            Inches(12.333), Inches(1.5)
        )
        msg_box.fill.solid()
        msg_box.fill.fore_color.rgb = hex_to_rgb(colors["darkblue"]["hex"])
        msg_box.line.fill.background()

        msg_text = slide.shapes.add_textbox(
            Inches(0.7), Inches(1.8),
            Inches(11.933), Inches(1)
        )
        msg_frame = msg_text.text_frame
        msg_para = msg_frame.paragraphs[0]
        msg_para.text = spec.content["main_message"]
        msg_para.font.size = Pt(24)
        msg_para.font.bold = True
        msg_para.font.color.rgb = RGBColor(255, 255, 255)
        msg_para.alignment = PP_ALIGN.CENTER

    # Action Items
    if spec.content and spec.content.get("action_items"):
        items = spec.content["action_items"]

        for i, item in enumerate(items[:5]):
            y = 3.5 + (i * 0.8)

            # Checkbox
            checkbox = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(1), Inches(y),
                Inches(0.4), Inches(0.4)
            )
            checkbox.fill.background()
            checkbox.line.color.rgb = hex_to_rgb(colors["darkblue"]["hex"])
            checkbox.line.width = Pt(2)

            # Text
            item_box = slide.shapes.add_textbox(
                Inches(1.6), Inches(y),
                Inches(10), Inches(0.6)
            )
            item_frame = item_box.text_frame
            item_para = item_frame.paragraphs[0]
            item_para.text = item
            item_para.font.size = Pt(18)
            item_para.font.color.rgb = hex_to_rgb(colors["darkgray"]["hex"])

    return slide


def add_sources_slide(prs: Presentation, spec: SlideSpec, style: Dict[str, Any]):
    """Fügt Quellen-Slide hinzu"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    colors = style["colors"]["primary"]
    typography = style["typography"]

    # Slide-Titel
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(12.333), Inches(1)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = spec.title or "Quellen & Kontakt"
    title_para.font.size = Pt(typography["sizes"]["slide_title"])
    title_para.font.bold = True
    title_para.font.color.rgb = hex_to_rgb(colors["darkblue"]["hex"])

    # Quellen-Liste
    if spec.content and spec.content.get("sources"):
        sources_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5),
            Inches(6), Inches(4)
        )
        sources_frame = sources_box.text_frame
        sources_frame.word_wrap = True

        for i, source in enumerate(spec.content["sources"]):
            if i == 0:
                para = sources_frame.paragraphs[0]
            else:
                para = sources_frame.add_paragraph()
            para.text = f"• {source}"
            para.font.size = Pt(14)
            para.font.color.rgb = hex_to_rgb(colors["darkgray"]["hex"])
            para.space_after = Pt(8)

    # Kontakt-Box
    if spec.content and spec.content.get("contact"):
        contact = spec.content["contact"]

        contact_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(7), Inches(1.5),
            Inches(5.833), Inches(3)
        )
        contact_box.fill.solid()
        contact_box.fill.fore_color.rgb = hex_to_rgb(colors["lightgray"]["hex"])
        contact_box.line.fill.background()

        # Company
        company_box = slide.shapes.add_textbox(
            Inches(7.3), Inches(1.8),
            Inches(5.233), Inches(0.6)
        )
        company_frame = company_box.text_frame
        company_para = company_frame.paragraphs[0]
        company_para.text = contact.get("company", "FehrAdvice & Partners AG")
        company_para.font.size = Pt(18)
        company_para.font.bold = True
        company_para.font.color.rgb = hex_to_rgb(colors["darkblue"]["hex"])

        # Details
        details_box = slide.shapes.add_textbox(
            Inches(7.3), Inches(2.5),
            Inches(5.233), Inches(1.5)
        )
        details_frame = details_box.text_frame
        details_frame.word_wrap = True

        details_para = details_frame.paragraphs[0]
        details_para.text = contact.get("email", "info@fehradvice.com")
        details_para.font.size = Pt(14)
        details_para.font.color.rgb = hex_to_rgb(colors["darkgray"]["hex"])

        web_para = details_frame.add_paragraph()
        web_para.text = contact.get("web", "www.fehradvice.com")
        web_para.font.size = Pt(14)
        web_para.font.color.rgb = hex_to_rgb(colors["darkblue"]["hex"])

    return slide


def add_generic_slide(prs: Presentation, spec: SlideSpec, style: Dict[str, Any]):
    """Fügt generischen Slide hinzu"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    colors = style["colors"]["primary"]
    typography = style["typography"]

    # Slide-Titel
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(12.333), Inches(1)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = spec.title or spec.slide_type.replace("_", " ").title()
    title_para.font.size = Pt(typography["sizes"]["slide_title"])
    title_para.font.bold = True
    title_para.font.color.rgb = hex_to_rgb(colors["darkblue"]["hex"])

    # Placeholder Content
    content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5),
        Inches(12.333), Inches(5)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    content_para = content_frame.paragraphs[0]
    content_para.text = f"[Inhalt für {spec.slide_type}]"
    content_para.font.size = Pt(18)
    content_para.font.color.rgb = hex_to_rgb(colors["darkgray"]["hex"])

    return slide


def generate_pptx(
    session_id: str,
    audience: str = "board",
    custom_8d: Optional[str] = None,
    output_path: Optional[Path] = None,
    use_master_template: bool = True
) -> Path:
    """Hauptfunktion: Generiert PPTX aus Session-Daten

    Args:
        session_id: EBF Session ID
        audience: Zielgruppen-Profil (board, management, team, science, client)
        custom_8d: Custom 8D-Profil als String
        output_path: Optionaler Output-Pfad
        use_master_template: Wenn True und FehrAdvice-Master.pptx existiert, wird es als Basis verwendet
    """

    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx ist nicht installiert. Bitte installieren mit: pip install python-pptx")

    # Lade Konfiguration
    style = load_style()
    mapping = load_8d_mapping()

    # Lade Session-Daten
    session_data = load_session_data(session_id)

    # Bestimme 8D-Profil
    if custom_8d:
        profile = parse_custom_8d(custom_8d)
    else:
        profile = get_audience_profile(audience, mapping)

    # Generiere Slide-Architektur
    slides = generate_slide_architecture(profile, session_data, mapping)

    # Erstelle Präsentation
    # Option 1: Verwende Master-Template (wenn vorhanden)
    # Option 2: Erstelle neue Präsentation
    if use_master_template and MASTER_TEMPLATE.exists():
        print(f"  Using master template: {MASTER_TEMPLATE.name}")
        prs = Presentation(str(MASTER_TEMPLATE))
        # Entferne Beispiel-Slides aus dem Master-Template
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
    else:
        prs = Presentation()

    prs.slide_width = Emu(style["slide"]["dimensions"]["width_emu"])
    prs.slide_height = Emu(style["slide"]["dimensions"]["height_emu"])

    # Füge Slides hinzu
    slide_generators = {
        "TITLE": add_title_slide,
        "EXECUTIVE_SUMMARY": add_executive_summary_slide,
        "DATA_CHART": lambda p, s, st: add_data_chart_slide(p, s, st, session_id),
        "IMPLICATIONS": add_implications_slide,
        "CALL_TO_ACTION": add_call_to_action_slide,
        "SOURCES": add_sources_slide,
    }

    for spec in slides:
        generator = slide_generators.get(spec.slide_type, add_generic_slide)
        if spec.slide_type == "DATA_CHART":
            generator(prs, spec, style)
        else:
            generator(prs, spec, style)

    # Speichere PPTX
    if output_path is None:
        session_dir = OUTPUTS_DIR / session_id
        session_dir.mkdir(parents=True, exist_ok=True)
        output_path = session_dir / f"presentation_{audience}.pptx"

    prs.save(str(output_path))

    print(f"✅ PPTX generiert: {output_path}")
    print(f"   Zielgruppe: {audience}")
    print(f"   Slides: {len(slides)}")

    return output_path


# =============================================================================
# CLI
# =============================================================================

def main():
    parser = argparse.ArgumentParser(
        description="EBF PPTX Generator - Generiert PowerPoint aus Session-Daten"
    )

    parser.add_argument(
        "--session", "-s",
        required=True,
        help="Session-ID (z.B. EBF-S-2026-01-26-COG-001)"
    )

    parser.add_argument(
        "--audience", "-a",
        default="board",
        choices=["board", "management", "team", "science", "client"],
        help="Zielgruppen-Profil (default: board)"
    )

    parser.add_argument(
        "--custom-8d", "-c",
        help="Custom 8D-Profil: 'D1,D2,D3,D4,D5,D6,D7,D8' (z.B. '0.4,0.3,0.85,0.2,G2,0.9,0.4,0.6')"
    )

    parser.add_argument(
        "--output", "-o",
        help="Output-Pfad (optional, default: outputs/sessions/<session>/presentation_<audience>.pptx)"
    )

    parser.add_argument(
        "--generate-graphics", "-g",
        action="store_true",
        help="Generiere zuerst die Grafiken mit generate_graphics.py"
    )

    parser.add_argument(
        "--list-audiences",
        action="store_true",
        help="Liste verfügbare Zielgruppen-Profile"
    )

    parser.add_argument(
        "--no-master",
        action="store_true",
        help="Erstelle Präsentation ohne Master-Template (nur programmatische Layouts)"
    )

    args = parser.parse_args()

    # Liste Audiences
    if args.list_audiences:
        mapping = load_8d_mapping()
        print("\nVerfügbare Zielgruppen-Profile:")
        print("=" * 60)
        for key, profile in mapping.get("audience_profiles", {}).items():
            print(f"\n{key}: {profile['name']}")
            print(f"  {profile['description']}")
            p = profile['profile_8d']
            print(f"  D1={p['D1_knowledge']}, D4={p['D4_time']}, D5={p['D5_goal']}")
        return

    # Generiere Grafiken falls gewünscht
    if args.generate_graphics:
        import subprocess
        print("Generiere Grafiken...")
        subprocess.run([
            "python", str(BASE_DIR / "scripts" / "generate_graphics.py"),
            "--session", args.session
        ])

    # Generiere PPTX
    output_path = Path(args.output) if args.output else None
    generate_pptx(
        session_id=args.session,
        audience=args.audience,
        custom_8d=args.custom_8d,
        output_path=output_path,
        use_master_template=not args.no_master
    )


if __name__ == "__main__":
    main()
