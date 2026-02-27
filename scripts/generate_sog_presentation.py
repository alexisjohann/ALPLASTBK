#!/usr/bin/env python3
"""
Generate SOG Presentation with FehrAdvice CI
=============================================

Converts "Ordnung im digitalen Raum" PDF to FehrAdvice Corporate Design.

FehrAdvice CI:
- Primary: Dunkelblau #024079, Hellblau #549EDE
- Text: Dunkelgrau #25212A, Hellgrau #F3F5F7
- Fonts: Roboto Bold (Headlines), Open Sans (Body)

Usage:
    python scripts/generate_sog_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Alias for consistency
RgbColor = RGBColor

# =============================================================================
# FEHRADVICE CORPORATE IDENTITY
# =============================================================================

class FehrAdviceCI:
    """FehrAdvice Corporate Identity Colors and Fonts"""

    # Primary Colors
    DARKBLUE = RgbColor(0x02, 0x40, 0x79)      # #024079 - Headlines, Akzente
    LIGHTBLUE = RgbColor(0x54, 0x9E, 0xDE)     # #549EDE - Sekundär
    DARKGRAY = RgbColor(0x25, 0x21, 0x2A)      # #25212A - Fliesstext
    LIGHTGRAY = RgbColor(0xF3, 0xF5, 0xF7)     # #F3F5F7 - Hintergründe

    # Secondary Colors
    MINT = RgbColor(0x7E, 0xBD, 0xAC)          # #7EBDAC - Positive
    OCHER = RgbColor(0xDE, 0xCB, 0x3F)         # #DECB3F - Warnung
    ORANGE = RgbColor(0xDE, 0x9D, 0x3E)        # #DE9D3E - Highlight

    # Neutral
    WHITE = RgbColor(0xFF, 0xFF, 0xFF)
    BLACK = RgbColor(0x00, 0x00, 0x00)

    # Fonts
    FONT_HEADLINE = "Roboto"
    FONT_BODY = "Open Sans"

    # Slide dimensions (16:9)
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)


def set_text_style(paragraph, font_name, font_size, color, bold=False, italic=False):
    """Apply consistent text styling"""
    paragraph.font.name = font_name
    paragraph.font.size = Pt(font_size)
    paragraph.font.color.rgb = color
    paragraph.font.bold = bold
    paragraph.font.italic = italic


def add_title_shape(slide, text, top=Inches(0.3), font_size=36):
    """Add a styled title to slide"""
    title_box = slide.shapes.add_textbox(
        Inches(0.5), top, Inches(12.333), Inches(1.0)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    set_text_style(p, FehrAdviceCI.FONT_HEADLINE, font_size, FehrAdviceCI.DARKBLUE, bold=True)
    return title_box


def add_subtitle(slide, text, top=Inches(1.1), font_size=24):
    """Add subtitle text"""
    box = slide.shapes.add_textbox(
        Inches(0.5), top, Inches(12.333), Inches(0.8)
    )
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    set_text_style(p, FehrAdviceCI.FONT_BODY, font_size, FehrAdviceCI.DARKGRAY)
    return box


def add_quote_box(slide, text, top, height=Inches(1.2)):
    """Add a highlighted quote box with FehrAdvice styling"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), top, Inches(12.333), height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = FehrAdviceCI.LIGHTGRAY
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    set_text_style(p, FehrAdviceCI.FONT_BODY, 18, FehrAdviceCI.DARKBLUE, bold=True, italic=True)
    return shape


def add_footer(slide, page_num, total_pages):
    """Add footer with page number"""
    footer = slide.shapes.add_textbox(
        Inches(11.5), Inches(6.9), Inches(1.5), Inches(0.4)
    )
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = f"{page_num} / {total_pages}"
    p.alignment = PP_ALIGN.RIGHT
    set_text_style(p, FehrAdviceCI.FONT_BODY, 11, FehrAdviceCI.DARKGRAY)


def add_status_badge(slide, text="Status: Strategisches Briefing"):
    """Add status badge in bottom right"""
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(10.0), Inches(6.7), Inches(2.8), Inches(0.5)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = FehrAdviceCI.DARKBLUE
    badge.line.fill.background()

    tf = badge.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    set_text_style(p, FehrAdviceCI.FONT_BODY, 12, FehrAdviceCI.WHITE, bold=True)


# =============================================================================
# SLIDE GENERATORS
# =============================================================================

def create_slide_01_title(prs):
    """Slide 1: Title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Main title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(2))
    tf = title.text_frame
    tf.word_wrap = True

    p1 = tf.paragraphs[0]
    p1.text = "Ordnung im digitalen Raum:"
    set_text_style(p1, FehrAdviceCI.FONT_HEADLINE, 44, FehrAdviceCI.DARKBLUE, bold=True)

    p2 = tf.add_paragraph()
    p2.text = "Das Social Media Ordnungs-Gesetz (SOG)"
    set_text_style(p2, FehrAdviceCI.FONT_HEADLINE, 44, FehrAdviceCI.DARKBLUE, bold=True)

    # Subtitle
    add_subtitle(slide, "Strategisches Briefing zur Kommunikationsoffensive.", Inches(2.8))

    # Visual: Dots pattern (simplified as description)
    desc = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12), Inches(1.5))
    tf = desc.text_frame
    p = tf.paragraphs[0]
    p.text = "[ Visualisierung: Chaos → Ordnung - Punkte die sich von chaotisch zu geordnet formieren ]"
    set_text_style(p, FehrAdviceCI.FONT_BODY, 14, FehrAdviceCI.DARKGRAY, italic=True)
    p.alignment = PP_ALIGN.CENTER

    add_status_badge(slide)
    return slide


def create_slide_02_summary(prs):
    """Slide 2: Management Summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_shape(slide, "Management Summary: Wir schaffen\nFairness und Verantwortung.", font_size=36)

    # Three columns: Situation - Komplikation - Lösung
    cols = [
        ("SITUATION", "📱", "Der digitale Raum ist im Vergleich zu analogen Medien (TV, Radio) weitgehend unreguliert."),
        ("KOMPLIKATION", "⚠️", "Dies gefährdet Kinder durch mangelnden Schutz und die Demokratie durch Desinformation und Polarisierung."),
        ("LÖSUNG", "⚖️", "Das SOG (Social Media Ordnungs-Gesetz). Kein Verbot, sondern die Einführung von fairen Spielregeln, wie sie offline längst gelten.")
    ]

    for i, (header, icon, text) in enumerate(cols):
        x = Inches(0.5 + i * 4.2)

        # Header
        hdr = slide.shapes.add_textbox(x, Inches(2.0), Inches(4), Inches(0.5))
        p = hdr.text_frame.paragraphs[0]
        p.text = header
        p.alignment = PP_ALIGN.CENTER
        set_text_style(p, FehrAdviceCI.FONT_HEADLINE, 18, FehrAdviceCI.DARKBLUE, bold=True)

        # Icon placeholder
        icn = slide.shapes.add_textbox(x, Inches(2.5), Inches(4), Inches(0.8))
        p = icn.text_frame.paragraphs[0]
        p.text = icon
        p.alignment = PP_ALIGN.CENTER
        set_text_style(p, FehrAdviceCI.FONT_BODY, 48, FehrAdviceCI.DARKBLUE)

        # Text
        txt = slide.shapes.add_textbox(x, Inches(3.4), Inches(3.8), Inches(2))
        tf = txt.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        set_text_style(p, FehrAdviceCI.FONT_BODY, 16, FehrAdviceCI.DARKGRAY)

    # Quote at bottom
    add_quote_box(slide, "«Diese Regierung schafft Ordnung. Wer Reichweite hat, trägt Verantwortung.»", Inches(5.8))

    add_footer(slide, 2, 12)
    return slide


def create_slide_03_frame(prs):
    """Slide 3: Strategic Frame - Von Verbot zu Ordnung"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1.2))
    tf = title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Der strategische Rahmen: Von «Verbot» zu «Ordnung»."
    set_text_style(p, FehrAdviceCI.FONT_HEADLINE, 32, FehrAdviceCI.DARKBLUE, bold=True)
    p2 = tf.add_paragraph()
    p2.text = "Wir verlassen die Defensive und argumentieren staatstragend."
    set_text_style(p2, FehrAdviceCI.FONT_BODY, 24, FehrAdviceCI.DARKGRAY, italic=True)

    # Two columns
    # Left: "Nicht mehr sagen" (strikethrough)
    left_hdr = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(5.5), Inches(0.6))
    p = left_hdr.text_frame.paragraphs[0]
    p.text = "NICHT MEHR SAGEN"
    set_text_style(p, FehrAdviceCI.FONT_HEADLINE, 18, FehrAdviceCI.DARKGRAY, bold=True)

    old_terms = ["Verbot / Einschränkung", "Ich / Mein Ministerium", "Strafen", "Defensive Erklärungen"]
    for i, term in enumerate(old_terms):
        box = slide.shapes.add_textbox(Inches(0.5), Inches(2.7 + i * 0.8), Inches(5), Inches(0.6))
        p = box.text_frame.paragraphs[0]
        p.text = f"—{term}—"  # Strikethrough effect
        set_text_style(p, FehrAdviceCI.FONT_BODY, 20, FehrAdviceCI.DARKGRAY)

    # Right: "Neu: Strategische Tonalität"
    right_hdr = slide.shapes.add_textbox(Inches(7), Inches(2.0), Inches(5.5), Inches(0.6))
    p = right_hdr.text_frame.paragraphs[0]
    p.text = "NEU: STRATEGISCHE TONALITÄT"
    set_text_style(p, FehrAdviceCI.FONT_HEADLINE, 18, FehrAdviceCI.DARKBLUE, bold=True)

    new_terms = ["Ordnung / Regeln / Schutz", "Diese Regierung / Wir als Regierung",
                 "Verantwortung einfordern / Konsequenzen", "Gleiche Spielregeln für alle"]
    for i, term in enumerate(new_terms):
        box = slide.shapes.add_textbox(Inches(7), Inches(2.7 + i * 0.8), Inches(5.5), Inches(0.6))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = f"→ {term}"
        set_text_style(p, FehrAdviceCI.FONT_BODY, 20, FehrAdviceCI.DARKBLUE, bold=True)

    # Bottom quote
    quote = slide.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(12.333), Inches(0.8))
    p = quote.text_frame.paragraphs[0]
    p.text = "Wir rechtfertigen uns nicht für Regeln. Wir schaffen sie."
    p.alignment = PP_ALIGN.CENTER
    set_text_style(p, FehrAdviceCI.FONT_HEADLINE, 24, FehrAdviceCI.DARKBLUE, bold=True)

    add_footer(slide, 3, 12)
    return slide


def create_slide_04_pillars(prs):
    """Slide 4: Four Pillars"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_shape(slide, "Vier Säulen für die digitale Souveränität.", font_size=36)

    pillars = [
        ("🛡️", "Altersgerechter\nZugang", "Schutz für Kinder und Jugendliche"),
        ("⚖️", "Plattform-\nHaftung", "Wer verdient, haftet"),
        ("📖", "Medien-\nkompetenz", "Demokratie braucht mündige Bürger:innen"),
        ("👁️", "Transparenz", "Gleiche Spielregeln für alle Medien")
    ]

    for i, (icon, title, desc) in enumerate(pillars):
        x = Inches(0.5 + i * 3.2)

        # Pillar header (blue background)
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x, Inches(1.8), Inches(2.9), Inches(1.5)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = FehrAdviceCI.DARKBLUE
        header.line.fill.background()

        # Icon
        icn = slide.shapes.add_textbox(x, Inches(1.9), Inches(2.9), Inches(0.8))
        p = icn.text_frame.paragraphs[0]
        p.text = icon
        p.alignment = PP_ALIGN.CENTER
        set_text_style(p, FehrAdviceCI.FONT_BODY, 36, FehrAdviceCI.WHITE)

        # Pillar body (light gray)
        body = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x, Inches(3.3), Inches(2.9), Inches(2.5)
        )
        body.fill.solid()
        body.fill.fore_color.rgb = FehrAdviceCI.LIGHTGRAY
        body.line.fill.background()

        # Title
        ttl = slide.shapes.add_textbox(x + Inches(0.1), Inches(3.5), Inches(2.7), Inches(1))
        tf = ttl.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        set_text_style(p, FehrAdviceCI.FONT_HEADLINE, 20, FehrAdviceCI.DARKBLUE, bold=True)

        # Description
        dsc = slide.shapes.add_textbox(x + Inches(0.1), Inches(4.6), Inches(2.7), Inches(1))
        tf = dsc.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.alignment = PP_ALIGN.CENTER
        set_text_style(p, FehrAdviceCI.FONT_BODY, 14, FehrAdviceCI.DARKGRAY)

    add_footer(slide, 4, 12)
    return slide


def create_slide_05_pillar1(prs):
    """Slide 5: Säule 1 - Kinderschutz"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_shape(slide, "Säule 1: Konsequenter Schutz für unsere Kinder.", font_size=32)

    # Quote box
    quote = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(4))
    tf = quote.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "«Unter 15 Jahren brauchen Kinder Schutz. Das ist konsequenter Jugendschutz – so wie bei Filmen, so wie bei Alkohol.»"
    set_text_style(p, FehrAdviceCI.FONT_HEADLINE, 28, FehrAdviceCI.DARKBLUE, bold=True, italic=True)

    # Icon description (right side)
    icon_box = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.5), Inches(3))
    tf = icon_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "15+ = Kein Alkohol unter 15 = Schutz auf Social Media"
    p.alignment = PP_ALIGN.CENTER
    set_text_style(p, FehrAdviceCI.FONT_BODY, 18, FehrAdviceCI.DARKGRAY)

    # Strategic standpoint box
    strat = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7), Inches(4.5), Inches(5.5), Inches(1.8)
    )
    strat.fill.solid()
    strat.fill.fore_color.rgb = FehrAdviceCI.LIGHTGRAY
    strat.line.fill.background()

    strat_txt = slide.shapes.add_textbox(Inches(7.2), Inches(4.6), Inches(5.1), Inches(1.6))
    tf = strat_txt.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Strategischer Standpunkt:"
    set_text_style(p, FehrAdviceCI.FONT_HEADLINE, 16, FehrAdviceCI.DARKBLUE, bold=True)
    p2 = tf.add_paragraph()
    p2.text = "Wenn Europa nicht schnell genug ist, gehen wir national voran. Das ist Ordnung schaffen. Nicht zusehen."
    set_text_style(p2, FehrAdviceCI.FONT_BODY, 14, FehrAdviceCI.DARKGRAY)

    add_footer(slide, 5, 12)
    return slide


def create_slide_06_pillar2(prs):
    """Slide 6: Säule 2 - Verantwortung und Haftung"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_shape(slide, "Säule 2: Verantwortung und Haftung.", font_size=36)

    # Subtitle
    sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.333), Inches(1))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.text = "Wer Reichweite hat, trägt Verantwortung.\nWer verdient, haftet."
    p.alignment = PP_ALIGN.CENTER
    set_text_style(p, FehrAdviceCI.FONT_HEADLINE, 24, FehrAdviceCI.DARKBLUE, bold=True)

    # Scale visualization description
    scale = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = scale.text_frame
    p = tf.paragraphs[0]
    p.text = "[ Visualisierung: Waage mit «Gewinn» vs «Haftung (DSA Rules)» ]"
    p.alignment = PP_ALIGN.CENTER
    set_text_style(p, FehrAdviceCI.FONT_BODY, 16, FehrAdviceCI.DARKGRAY, italic=True)

    # Two columns
    # Left: Der Mechanismus
    left = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(5.5), Inches(1.5))
    tf = left.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Der Mechanismus"
    set_text_style(p, FehrAdviceCI.FONT_HEADLINE, 20, FehrAdviceCI.DARKBLUE, bold=True)
    p2 = tf.add_paragraph()
    p2.text = "Der DSA (Digital Services Act) gibt den Rahmen. Strafen bis zu 6% des Jahresumsatzes."
    set_text_style(p2, FehrAdviceCI.FONT_BODY, 16, FehrAdviceCI.DARKGRAY)

    # Right: Das Framing
    right = slide.shapes.add_textbox(Inches(7), Inches(4.5), Inches(5.5), Inches(1.5))
    tf = right.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Das Framing"
    set_text_style(p, FehrAdviceCI.FONT_HEADLINE, 20, FehrAdviceCI.DARKBLUE, bold=True)
    p2 = tf.add_paragraph()
    p2.text = "Im Fernsehen und Radio gelten Regeln. Das muss auch für Social Media gelten. Gleiche Spielregeln für alle Medien sind keine Einschränkung – das ist Fairness."
    set_text_style(p2, FehrAdviceCI.FONT_BODY, 16, FehrAdviceCI.DARKGRAY)

    add_footer(slide, 6, 12)
    return slide


def create_slide_07_pillar3(prs):
    """Slide 7: Säule 3 - Medienkompetenz"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_shape(slide, "Säule 3: Medienkompetenz ist Demokratiepolitik.", font_size=32)

    # Two boxes top
    # Left: Kontext
    ctx = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.5), Inches(5.8), Inches(1.5)
    )
    ctx.fill.solid()
    ctx.fill.fore_color.rgb = FehrAdviceCI.LIGHTGRAY
    ctx.line.fill.background()

    ctx_txt = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(5.4), Inches(1.3))
    tf = ctx_txt.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Kontext: Junge Menschen informieren sich über Plattformen. Dort verschwimmen Journalismus und Desinformation."
    set_text_style(p, FehrAdviceCI.FONT_BODY, 16, FehrAdviceCI.DARKGRAY)

    # Right: Ziel
    goal = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.8), Inches(1.5), Inches(5.8), Inches(1.5)
    )
    goal.fill.solid()
    goal.fill.fore_color.rgb = FehrAdviceCI.LIGHTGRAY
    goal.line.fill.background()

    goal_txt = slide.shapes.add_textbox(Inches(7), Inches(1.6), Inches(5.4), Inches(1.3))
    tf = goal_txt.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Ziel: Demokratie braucht mündige Bürger:innen, die Quellen bewerten können."
    set_text_style(p, FehrAdviceCI.FONT_HEADLINE, 18, FehrAdviceCI.DARKBLUE, bold=True)

    # Statement
    stmt = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(8), Inches(1))
    tf = stmt.text_frame
    p = tf.paragraphs[0]
    p.text = "Diese Regierung wird nicht nur darüber reden.\nWir bauen die Struktur auf."
    set_text_style(p, FehrAdviceCI.FONT_HEADLINE, 24, FehrAdviceCI.DARKBLUE, bold=True)

    # Timeline indicator
    timeline = slide.shapes.add_textbox(Inches(9), Inches(4.5), Inches(3.5), Inches(1.5))
    tf = timeline.text_frame
    p = tf.paragraphs[0]
    p.text = "Start der Massnahmen:\n2026."
    p.alignment = PP_ALIGN.CENTER
    set_text_style(p, FehrAdviceCI.FONT_HEADLINE, 24, FehrAdviceCI.DARKBLUE, bold=True)

    add_footer(slide, 7, 12)
    return slide


def create_slide_08_pillar4(prs):
    """Slide 8: Säule 4 - Transparenz"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_shape(slide, "Säule 4: Transparenz und Algorithmen.", font_size=36)
    add_subtitle(slide, "Keine Sonderrechte für Big Tech.", Inches(1.0), 24)

    # Comparison table
    # Left: Traditionelle Medien
    left_hdr = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(5.5), Inches(0.8))
    tf = left_hdr.text_frame
    p = tf.paragraphs[0]
    p.text = "Traditionelle Medien\n(TV/Print)"
    set_text_style(p, FehrAdviceCI.FONT_HEADLINE, 20, FehrAdviceCI.DARKBLUE, bold=True)

    left_txt = slide.shapes.add_textbox(Inches(0.5), Inches(2.9), Inches(5.5), Inches(1.2))
    tf = left_txt.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Fake News verboten.\nBetrügerische Werbung verboten. ✓"
    set_text_style(p, FehrAdviceCI.FONT_BODY, 18, FehrAdviceCI.DARKGRAY)

    # Right: Social Media
    right_hdr = slide.shapes.add_textbox(Inches(7), Inches(2.0), Inches(5.5), Inches(0.8))
    tf = right_hdr.text_frame
    p = tf.paragraphs[0]
    p.text = "Social Media\nPlattformen"
    set_text_style(p, FehrAdviceCI.FONT_HEADLINE, 20, FehrAdviceCI.DARKBLUE, bold=True)

    right_txt = slide.shapes.add_textbox(Inches(7), Inches(2.9), Inches(5.5), Inches(1.2))
    tf = right_txt.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Verdienen Milliarden mit Polarisierung, Radikalisierung und Suchtverhalten. ⚠️"
    set_text_style(p, FehrAdviceCI.FONT_BODY, 18, FehrAdviceCI.DARKGRAY)

    # Bottom quote box
    quote_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(5.0), Inches(12.333), Inches(1.2)
    )
    quote_box.fill.solid()
    quote_box.fill.fore_color.rgb = FehrAdviceCI.DARKBLUE
    quote_box.line.fill.background()

    quote = slide.shapes.add_textbox(Inches(0.7), Inches(5.2), Inches(12), Inches(0.8))
    tf = quote.text_frame
    p = tf.paragraphs[0]
    p.text = "All das ist klassischen Medien verboten. Warum sollten wir es Social Media erlauben? Gleiche Spielregeln für alle."
    p.alignment = PP_ALIGN.CENTER
    set_text_style(p, FehrAdviceCI.FONT_HEADLINE, 18, FehrAdviceCI.WHITE, bold=True)

    add_footer(slide, 8, 12)
    return slide


def create_slide_09_counter(prs):
    """Slide 9: Counter-Narrativ FPÖ"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title (italic style)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Counter-Narrativ: Die «Zensur»-Lüge der FPÖ."
    set_text_style(p, FehrAdviceCI.FONT_HEADLINE, 32, FehrAdviceCI.DARKBLUE, bold=True, italic=True)

    # Speech bubble (left)
    bubble = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.5), Inches(4.5), Inches(2.5)
    )
    bubble.fill.solid()
    bubble.fill.fore_color.rgb = FehrAdviceCI.LIGHTGRAY
    bubble.line.color.rgb = FehrAdviceCI.DARKGRAY

    bubble_txt = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(4), Inches(1.5))
    tf = bubble_txt.text_frame
    p = tf.paragraphs[0]
    p.text = "Die FPÖ ruft:\nZensur!"
    p.alignment = PP_ALIGN.CENTER
    set_text_style(p, FehrAdviceCI.FONT_HEADLINE, 28, FehrAdviceCI.DARKBLUE, bold=True)

    # Strategic answer (right)
    ans_hdr = slide.shapes.add_textbox(Inches(6), Inches(1.3), Inches(6.5), Inches(0.6))
    tf = ans_hdr.text_frame
    p = tf.paragraphs[0]
    p.text = "Die strategische Antwort"
    set_text_style(p, FehrAdviceCI.FONT_HEADLINE, 22, FehrAdviceCI.DARKBLUE, bold=True)

    answers = [
        "Ich erinnere daran: Unter FPÖ-Regierungsbeteiligung ist in Sachen Plattformregulierung nichts passiert. Null.",
        "Die rufen ‚Zensur' – und meinen: Keine Regeln für ihre Freunde in den Tech-Konzernen.",
        "Im Fernsehen gelten Regeln. Im Radio gelten Regeln. Warum soll ausgerechnet TikTok keine Regeln haben?"
    ]

    for i, ans in enumerate(answers):
        box = slide.shapes.add_textbox(Inches(6), Inches(2.0 + i * 1.1), Inches(6.5), Inches(1))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"● {ans}"
        set_text_style(p, FehrAdviceCI.FONT_BODY, 14, FehrAdviceCI.DARKGRAY)

    # Bottom statement
    stmt = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12.333), Inches(0.8))
    tf = stmt.text_frame
    p = tf.paragraphs[0]
    p.text = "Das ist keine Zensur. Das ist Ordnung."
    p.alignment = PP_ALIGN.CENTER
    set_text_style(p, FehrAdviceCI.FONT_HEADLINE, 28, FehrAdviceCI.DARKBLUE, bold=True)

    add_footer(slide, 9, 12)
    return slide


def create_slide_10_faq(prs):
    """Slide 10: Critical Questions"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_shape(slide, "Antworten auf kritische Fragen der Umsetzung.", font_size=32)

    # Question 1
    q1_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.5), Inches(12.333), Inches(2.2)
    )
    q1_box.fill.solid()
    q1_box.fill.fore_color.rgb = FehrAdviceCI.LIGHTGRAY
    q1_box.line.fill.background()

    q1_txt = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.8), Inches(2))
    tf = q1_txt.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Frage 1: Werden Kinder das Verbot nicht umgehen (VPN)?"
    set_text_style(p, FehrAdviceCI.FONT_HEADLINE, 20, FehrAdviceCI.DARKBLUE, bold=True)
    p2 = tf.add_paragraph()
    p2.text = "«Natürlich werden manche es versuchen. Wie beim Alkohol. Aber die klare Grenze wirkt – sie schafft Bewusstsein. Wer weiss, dass etwas potenziell schädlich ist, geht anders damit um.»"
    set_text_style(p2, FehrAdviceCI.FONT_BODY, 16, FehrAdviceCI.DARKGRAY, italic=True)

    # Question 2
    q2_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(4.0), Inches(12.333), Inches(2.2)
    )
    q2_box.fill.solid()
    q2_box.fill.fore_color.rgb = FehrAdviceCI.LIGHTGRAY
    q2_box.line.fill.background()

    q2_txt = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(11.8), Inches(2))
    tf = q2_txt.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Frage 2: Sehen Experten ein Grundrechtsproblem (Datenschutz)?"
    set_text_style(p, FehrAdviceCI.FONT_HEADLINE, 20, FehrAdviceCI.DARKBLUE, bold=True)
    p2 = tf.add_paragraph()
    p2.text = "«Ein Eingriff ist zulässig, wenn er verhältnismässig ist. Alle Erkenntnisse zeigen: Der Schutzbedarf ist erheblich. Wir arbeiten daran, dass der Eingriff so gering wie möglich ist – aber wirksam bleibt.»"
    set_text_style(p2, FehrAdviceCI.FONT_BODY, 16, FehrAdviceCI.DARKGRAY, italic=True)

    add_footer(slide, 10, 12)
    return slide


def create_slide_11_timeline(prs):
    """Slide 11: Timeline"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1.5))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Wir warten nicht.\nWir handeln."
    set_text_style(p, FehrAdviceCI.FONT_HEADLINE, 40, FehrAdviceCI.DARKBLUE, bold=True)

    # Strategic ambition box
    amb_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7), Inches(1.5), Inches(5.5), Inches(1.5)
    )
    amb_box.fill.solid()
    amb_box.fill.fore_color.rgb = FehrAdviceCI.LIGHTGRAY
    amb_box.line.color.rgb = FehrAdviceCI.DARKBLUE

    amb_txt = slide.shapes.add_textbox(Inches(7.2), Inches(1.6), Inches(5.1), Inches(1.3))
    tf = amb_txt.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Strategische Ambition"
    set_text_style(p, FehrAdviceCI.FONT_HEADLINE, 18, FehrAdviceCI.DARKBLUE, bold=True)
    p2 = tf.add_paragraph()
    p2.text = "Wenn Europa nicht schnell genug ist, gehen wir national voran."
    set_text_style(p2, FehrAdviceCI.FONT_BODY, 14, FehrAdviceCI.DARKGRAY)

    # Timeline arrow
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(0.5), Inches(4.5), Inches(12.333), Inches(0.5)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = FehrAdviceCI.DARKBLUE
    arrow.line.fill.background()

    # Milestone 1: Bis Sommer
    m1 = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(5), Inches(1.5))
    tf = m1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Bis Sommer"
    set_text_style(p, FehrAdviceCI.FONT_HEADLINE, 20, FehrAdviceCI.DARKBLUE, bold=True)
    p2 = tf.add_paragraph()
    p2.text = "Vorlage des Gesetzesentwurfs (SOG). Diese Regierung schafft Ordnung im digitalen Raum."
    set_text_style(p2, FehrAdviceCI.FONT_BODY, 14, FehrAdviceCI.DARKGRAY)

    # Milestone 2: 2026
    m2 = slide.shapes.add_textbox(Inches(8), Inches(5.2), Inches(4.5), Inches(1.5))
    tf = m2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "2026"
    set_text_style(p, FehrAdviceCI.FONT_HEADLINE, 20, FehrAdviceCI.DARKBLUE, bold=True)
    p2 = tf.add_paragraph()
    p2.text = "Implementierung der Massnahmen zur Medienkompetenz."
    set_text_style(p2, FehrAdviceCI.FONT_BODY, 14, FehrAdviceCI.DARKGRAY)

    add_footer(slide, 11, 12)
    return slide


def create_slide_12_messages(prs):
    """Slide 12: 5 Key Messages"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1.2))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Kommunikation auf einen Blick:\nDie 5 Kernbotschaften."
    set_text_style(p, FehrAdviceCI.FONT_HEADLINE, 32, FehrAdviceCI.DARKBLUE, bold=True)

    sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.333), Inches(0.5))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.text = "Cheat Sheet für Interviews und Pressetermine."
    set_text_style(p, FehrAdviceCI.FONT_BODY, 18, FehrAdviceCI.DARKGRAY)

    messages = [
        ("1.", "SOG Allgemein", "«Diese Regierung schafft Ordnung im digitalen Raum.»"),
        ("2.", "Altersschutz", "«Kinder schützen. So wie bei Filmen, so bei Social Media.»"),
        ("3.", "Plattformen", "«Wer Reichweite hat, trägt Verantwortung.»"),
        ("4.", "Gegen FPÖ", "«Gleiche Spielregeln für alle Medien. Das ist keine Zensur – das ist Ordnung.»"),
        ("5.", "Geschwindigkeit", "«Wenn Europa nicht schnell genug ist, gehen wir voran.»")
    ]

    # Grid: 3 top, 2 bottom
    positions = [
        (Inches(0.5), Inches(2.2)),
        (Inches(4.7), Inches(2.2)),
        (Inches(8.9), Inches(2.2)),
        (Inches(0.5), Inches(4.5)),
        (Inches(4.7), Inches(4.5))
    ]

    for i, (num, title_text, quote) in enumerate(messages):
        x, y = positions[i]

        # Box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y, Inches(3.9), Inches(2.0)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = FehrAdviceCI.LIGHTGRAY
        box.line.color.rgb = FehrAdviceCI.DARKGRAY

        # Number
        num_txt = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.1), Inches(0.8), Inches(0.6))
        p = num_txt.text_frame.paragraphs[0]
        p.text = num
        set_text_style(p, FehrAdviceCI.FONT_HEADLINE, 28, FehrAdviceCI.DARKBLUE, bold=True)

        # Title
        ttl = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.6), Inches(3.7), Inches(0.5))
        p = ttl.text_frame.paragraphs[0]
        p.text = title_text
        set_text_style(p, FehrAdviceCI.FONT_HEADLINE, 16, FehrAdviceCI.DARKBLUE, bold=True)

        # Quote
        qt = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(1.1), Inches(3.7), Inches(0.8))
        tf = qt.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = quote
        set_text_style(p, FehrAdviceCI.FONT_BODY, 12, FehrAdviceCI.DARKGRAY)

    add_footer(slide, 12, 12)
    return slide


# =============================================================================
# MAIN
# =============================================================================

def main():
    """Generate the presentation"""
    print("=" * 60)
    print("Generating SOG Presentation with FehrAdvice CI")
    print("=" * 60)

    # Create presentation
    prs = Presentation()
    prs.slide_width = FehrAdviceCI.SLIDE_WIDTH
    prs.slide_height = FehrAdviceCI.SLIDE_HEIGHT

    # Generate all slides
    print("\nGenerating slides...")
    create_slide_01_title(prs)
    print("  ✓ Slide 1: Title")

    create_slide_02_summary(prs)
    print("  ✓ Slide 2: Management Summary")

    create_slide_03_frame(prs)
    print("  ✓ Slide 3: Strategic Frame")

    create_slide_04_pillars(prs)
    print("  ✓ Slide 4: Four Pillars")

    create_slide_05_pillar1(prs)
    print("  ✓ Slide 5: Säule 1 - Kinderschutz")

    create_slide_06_pillar2(prs)
    print("  ✓ Slide 6: Säule 2 - Haftung")

    create_slide_07_pillar3(prs)
    print("  ✓ Slide 7: Säule 3 - Medienkompetenz")

    create_slide_08_pillar4(prs)
    print("  ✓ Slide 8: Säule 4 - Transparenz")

    create_slide_09_counter(prs)
    print("  ✓ Slide 9: Counter-Narrativ")

    create_slide_10_faq(prs)
    print("  ✓ Slide 10: FAQ")

    create_slide_11_timeline(prs)
    print("  ✓ Slide 11: Timeline")

    create_slide_12_messages(prs)
    print("  ✓ Slide 12: Key Messages")

    # Save
    output_dir = "outputs/presentations"
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "SOG_FehrAdvice_CI.pptx")

    prs.save(output_path)

    print("\n" + "=" * 60)
    print(f"✅ PRESENTATION GENERATED: {output_path}")
    print("=" * 60)
    print(f"""
FehrAdvice CI Applied:
  • Colors: Dunkelblau #024079, Hellblau #549EDE
  • Fonts: Roboto Bold (Headlines), Open Sans (Body)
  • Layout: 16:9 format, consistent margins
  • Slides: 12 (1-to-1 from original PDF)

Next steps:
  1. Open {output_path} in PowerPoint
  2. Add FehrAdvice logo (templates/FAP_Logo.png)
  3. Adjust graphics/icons as needed
  4. Review and finalize
""")

    return output_path


if __name__ == "__main__":
    main()
