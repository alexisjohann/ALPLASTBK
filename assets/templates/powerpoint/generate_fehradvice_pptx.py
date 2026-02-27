#!/usr/bin/env python3
"""
FehrAdvice & Partners AG - PowerPoint Template Generator

Generates a branded PowerPoint presentation template following
the official Corporate Identity Guide.

Colors:
- Primary: Dark Blue (#024079), Light Blue (#549EDE)
- Neutral: Dark Gray (#25212A), Light Gray (#F3F5F7)
- Secondary: Lilac (#A1A0C6), Mint (#7EBDAC), Ocher (#DECB3F), Orange (#DE9D3E)

Fonts: Roboto (headlines), Open Sans (body)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os


class RgbColor:
    """Simple RGB color wrapper."""
    def __init__(self, r, g, b):
        self.r = r
        self.g = g
        self.b = b

    def __iter__(self):
        return iter((self.r, self.g, self.b))

# =============================================================================
# FehrAdvice Brand Colors (from CI Guide)
# =============================================================================

COLORS = {
    # Primary
    'dark_blue': RgbColor(0x02, 0x40, 0x79),      # #024079
    'light_blue': RgbColor(0x54, 0x9E, 0xDE),     # #549EDE
    'dark_gray': RgbColor(0x25, 0x21, 0x2A),      # #25212A
    'light_gray': RgbColor(0xF3, 0xF5, 0xF7),     # #F3F5F7
    'white': RgbColor(0xFF, 0xFF, 0xFF),          # #FFFFFF

    # Secondary
    'lilac': RgbColor(0xA1, 0xA0, 0xC6),          # #A1A0C6
    'mint': RgbColor(0x7E, 0xBD, 0xAC),           # #7EBDAC
    'ocher': RgbColor(0xDE, 0xCB, 0x3F),          # #DECB3F
    'orange': RgbColor(0xDE, 0x9D, 0x3E),         # #DE9D3E
}

# Fonts
FONT_HEADLINE = 'Roboto'
FONT_BODY = 'Open Sans'
FONT_ACCENT = 'Playfair Display'

# Slide dimensions (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_font(run, name=FONT_BODY, size=18, bold=False, italic=False, color=None):
    """Apply font settings to a text run."""
    from pptx.dml.color import RGBColor
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = RGBColor(color.r, color.g, color.b)


def add_rectangle(slide, left, top, width, height, fill_color, line_color=None):
    """Add a colored rectangle shape."""
    from pptx.dml.color import RGBColor
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(fill_color.r, fill_color.g, fill_color.b)
    if line_color:
        shape.line.color.rgb = RGBColor(line_color.r, line_color.g, line_color.b)
    else:
        shape.line.fill.background()
    return shape


def create_title_slide(prs):
    """Slide 1: Title slide with dark blue background."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Dark blue background
    add_rectangle(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, COLORS['dark_blue'])

    # Light blue accent bar (left)
    add_rectangle(slide, 0, Inches(2.5), Inches(0.15), Inches(2), COLORS['light_blue'])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(10), Inches(1))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Präsentationstitel"
    set_font(run, FONT_HEADLINE, 44, bold=True, color=COLORS['white'])

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(10), Inches(0.8))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Untertitel der Präsentation"
    set_font(run, FONT_BODY, 24, color=COLORS['light_blue'])

    # Date/Author line
    info_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(10), Inches(0.5))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Januar 2026 | FehrAdvice & Partners AG"
    set_font(run, FONT_BODY, 14, color=COLORS['light_gray'])

    # Logo placeholder (bottom right)
    logo_box = slide.shapes.add_textbox(Inches(10.5), Inches(6.3), Inches(2.5), Inches(0.8))
    tf = logo_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "[LOGO]"
    set_font(run, FONT_BODY, 12, color=COLORS['light_gray'])

    return slide


def create_agenda_slide(prs):
    """Slide 2: Agenda slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Light gray background
    add_rectangle(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, COLORS['light_gray'])

    # Dark blue header bar
    add_rectangle(slide, 0, 0, SLIDE_WIDTH, Inches(1.2), COLORS['dark_blue'])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(10), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Agenda"
    set_font(run, FONT_HEADLINE, 32, bold=True, color=COLORS['white'])

    # Agenda items
    agenda_items = [
        ("01", "Ausgangslage und Ziele"),
        ("02", "Analyse und Erkenntnisse"),
        ("03", "Empfehlungen"),
        ("04", "Nächste Schritte"),
    ]

    y_pos = Inches(1.8)
    for num, text in agenda_items:
        # Number
        num_box = slide.shapes.add_textbox(Inches(0.8), y_pos, Inches(0.8), Inches(0.6))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = num
        set_font(run, FONT_HEADLINE, 28, bold=True, color=COLORS['light_blue'])

        # Text
        text_box = slide.shapes.add_textbox(Inches(1.8), y_pos, Inches(10), Inches(0.6))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        set_font(run, FONT_BODY, 22, color=COLORS['dark_gray'])

        y_pos += Inches(1.1)

    return slide


def create_section_slide(prs, section_num="01", section_title="Kapitelüberschrift"):
    """Section divider slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Dark blue background
    add_rectangle(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, COLORS['dark_blue'])

    # Light blue accent circle
    from pptx.dml.color import RGBColor
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.5), Inches(2.8), Inches(1.5), Inches(1.5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(COLORS['light_blue'].r, COLORS['light_blue'].g, COLORS['light_blue'].b)
    circle.line.fill.background()

    # Section number in circle
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(1.5), Inches(0.8))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = section_num
    set_font(run, FONT_HEADLINE, 36, bold=True, color=COLORS['white'])

    # Section title
    title_box = slide.shapes.add_textbox(Inches(2.5), Inches(3.0), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = section_title
    set_font(run, FONT_HEADLINE, 40, bold=True, color=COLORS['white'])

    return slide


def create_content_slide(prs, title="Inhaltsfolie", bullets=None):
    """Standard content slide with title and bullet points."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    if bullets is None:
        bullets = [
            "Erster Aufzählungspunkt mit wichtiger Information",
            "Zweiter Punkt erklärt weitere Details",
            "Dritter Punkt fasst Erkenntnisse zusammen",
            "Vierter Punkt gibt Handlungsempfehlung",
        ]

    # White background
    add_rectangle(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, COLORS['white'])

    # Dark blue header bar
    add_rectangle(slide, 0, 0, SLIDE_WIDTH, Inches(1.2), COLORS['dark_blue'])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(11), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_font(run, FONT_HEADLINE, 28, bold=True, color=COLORS['white'])

    # Light blue accent line
    add_rectangle(slide, Inches(0.5), Inches(1.5), Inches(2), Inches(0.08), COLORS['light_blue'])

    # Bullet points
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(12), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.level = 0
        p.space_before = Pt(12)
        p.space_after = Pt(6)

        # Add bullet character manually
        run = p.add_run()
        run.text = "●  " + bullet
        set_font(run, FONT_BODY, 18, color=COLORS['dark_gray'])

    # Footer
    add_slide_footer(slide)

    return slide


def create_two_column_slide(prs, title="Zwei-Spalten-Layout"):
    """Two-column content slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # White background
    add_rectangle(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, COLORS['white'])

    # Dark blue header bar
    add_rectangle(slide, 0, 0, SLIDE_WIDTH, Inches(1.2), COLORS['dark_blue'])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(11), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_font(run, FONT_HEADLINE, 28, bold=True, color=COLORS['white'])

    # Left column header
    left_header = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(0.5))
    tf = left_header.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Linke Spalte"
    set_font(run, FONT_HEADLINE, 20, bold=True, color=COLORS['light_blue'])

    # Left column content
    left_content = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(5.5), Inches(4))
    tf = left_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Inhalt der linken Spalte. Hier können Texte, Aufzählungen oder andere Inhalte platziert werden."
    set_font(run, FONT_BODY, 16, color=COLORS['dark_gray'])

    # Vertical divider
    add_rectangle(slide, Inches(6.4), Inches(1.5), Inches(0.03), Inches(5), COLORS['light_gray'])

    # Right column header
    right_header = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.5), Inches(0.5))
    tf = right_header.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Rechte Spalte"
    set_font(run, FONT_HEADLINE, 20, bold=True, color=COLORS['light_blue'])

    # Right column content
    right_content = slide.shapes.add_textbox(Inches(6.8), Inches(2.1), Inches(5.5), Inches(4))
    tf = right_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Inhalt der rechten Spalte. Ideal für Vergleiche, Gegenüberstellungen oder ergänzende Informationen."
    set_font(run, FONT_BODY, 16, color=COLORS['dark_gray'])

    # Footer
    add_slide_footer(slide)

    return slide


def create_quote_slide(prs, quote="Zitat hier einfügen", author="Autor"):
    """Quote/highlight slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Light blue background
    add_rectangle(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, COLORS['light_blue'])

    # Large quote mark
    quote_mark = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(2), Inches(1.5))
    tf = quote_mark.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "«"
    set_font(run, FONT_ACCENT, 120, color=COLORS['dark_blue'])

    # Quote text
    quote_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10), Inches(2.5))
    tf = quote_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = quote
    set_font(run, FONT_ACCENT, 32, italic=True, color=COLORS['white'])

    # Author
    author_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.2), Inches(10), Inches(0.5))
    tf = author_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"— {author}"
    set_font(run, FONT_BODY, 18, bold=True, color=COLORS['dark_blue'])

    return slide


def create_key_insight_slide(prs, insight="Kernaussage", details=None):
    """Key insight/takeaway slide with large text."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    if details is None:
        details = [
            "Unterstützender Punkt 1",
            "Unterstützender Punkt 2",
            "Unterstützender Punkt 3",
        ]

    # White background
    add_rectangle(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, COLORS['white'])

    # Dark blue left bar
    add_rectangle(slide, 0, 0, Inches(0.3), SLIDE_HEIGHT, COLORS['dark_blue'])

    # "KEY INSIGHT" label
    label_box = slide.shapes.add_textbox(Inches(0.8), Inches(1), Inches(3), Inches(0.4))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "KEY INSIGHT"
    set_font(run, FONT_HEADLINE, 14, bold=True, color=COLORS['light_blue'])

    # Main insight text
    insight_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11), Inches(2))
    tf = insight_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = insight
    set_font(run, FONT_HEADLINE, 36, bold=True, color=COLORS['dark_blue'])

    # Supporting points
    details_box = slide.shapes.add_textbox(Inches(0.8), Inches(4), Inches(11), Inches(2.5))
    tf = details_box.text_frame
    tf.word_wrap = True

    for i, detail in enumerate(details):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(8)
        run = p.add_run()
        run.text = f"→  {detail}"
        set_font(run, FONT_BODY, 18, color=COLORS['dark_gray'])

    # Footer
    add_slide_footer(slide)

    return slide


def create_chart_placeholder_slide(prs, title="Daten und Visualisierung"):
    """Slide with chart placeholder area."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # White background
    add_rectangle(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, COLORS['white'])

    # Dark blue header bar
    add_rectangle(slide, 0, 0, SLIDE_WIDTH, Inches(1.2), COLORS['dark_blue'])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(11), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_font(run, FONT_HEADLINE, 28, bold=True, color=COLORS['white'])

    # Chart placeholder area
    chart_area = add_rectangle(
        slide, Inches(0.5), Inches(1.5), Inches(8), Inches(5.2),
        COLORS['light_gray'], COLORS['light_blue']
    )

    # Placeholder text
    placeholder_box = slide.shapes.add_textbox(Inches(2.5), Inches(3.5), Inches(4), Inches(1))
    tf = placeholder_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "[Chart / Grafik hier einfügen]"
    set_font(run, FONT_BODY, 16, color=COLORS['dark_gray'])

    # Side notes area
    notes_header = slide.shapes.add_textbox(Inches(9), Inches(1.5), Inches(3.5), Inches(0.5))
    tf = notes_header.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Kernaussagen"
    set_font(run, FONT_HEADLINE, 16, bold=True, color=COLORS['light_blue'])

    notes = [
        "Erste Beobachtung aus den Daten",
        "Zweite wichtige Erkenntnis",
        "Implikation für die Praxis",
    ]

    notes_box = slide.shapes.add_textbox(Inches(9), Inches(2.1), Inches(3.5), Inches(4))
    tf = notes_box.text_frame
    tf.word_wrap = True

    for i, note in enumerate(notes):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(12)
        run = p.add_run()
        run.text = f"●  {note}"
        set_font(run, FONT_BODY, 14, color=COLORS['dark_gray'])

    # Footer
    add_slide_footer(slide)

    return slide


def create_contact_slide(prs):
    """Final contact/thank you slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Dark blue background
    add_rectangle(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, COLORS['dark_blue'])

    # Thank you text
    thanks_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(6), Inches(1))
    tf = thanks_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Vielen Dank"
    set_font(run, FONT_HEADLINE, 48, bold=True, color=COLORS['white'])

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(6), Inches(0.6))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "für Ihre Aufmerksamkeit"
    set_font(run, FONT_BODY, 24, color=COLORS['light_blue'])

    # Contact info box
    contact_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(5), Inches(2))
    tf = contact_box.text_frame
    tf.word_wrap = True

    contact_lines = [
        "FehrAdvice & Partners AG",
        "Klausstrasse 20",
        "8008 Zürich",
        "",
        "www.fehradvice.com",
        "info@fehradvice.com",
    ]

    for i, line in enumerate(contact_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        if i == 0:
            set_font(run, FONT_HEADLINE, 16, bold=True, color=COLORS['white'])
        elif line.startswith("www") or line.startswith("info"):
            set_font(run, FONT_BODY, 14, color=COLORS['light_blue'])
        else:
            set_font(run, FONT_BODY, 14, color=COLORS['light_gray'])

    # Light blue accent rectangle on right
    add_rectangle(slide, Inches(9), Inches(1.5), Inches(3.5), Inches(4.5), COLORS['light_blue'])

    # Logo placeholder in accent area
    logo_box = slide.shapes.add_textbox(Inches(9.5), Inches(3.2), Inches(2.5), Inches(1))
    tf = logo_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "[LOGO]"
    set_font(run, FONT_BODY, 16, color=COLORS['dark_blue'])

    return slide


def add_slide_footer(slide, page_num=None):
    """Add standard footer to slide."""
    # Company name
    company_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(3), Inches(0.4))
    tf = company_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "FehrAdvice & Partners AG"
    set_font(run, FONT_BODY, 10, color=COLORS['dark_gray'])

    # Confidential marker
    conf_box = slide.shapes.add_textbox(Inches(5), Inches(6.9), Inches(3), Inches(0.4))
    tf = conf_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Vertraulich"
    set_font(run, FONT_BODY, 10, color=COLORS['dark_gray'])


def create_fehradvice_template(output_path="FehrAdvice_Template.pptx"):
    """Generate complete FehrAdvice PowerPoint template."""

    # Create presentation with 16:9 aspect ratio
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Create all slide types
    create_title_slide(prs)
    create_agenda_slide(prs)
    create_section_slide(prs, "01", "Ausgangslage und Ziele")
    create_content_slide(prs, "Kontext und Herausforderung")
    create_two_column_slide(prs, "Vergleich: Ist vs. Soll")
    create_section_slide(prs, "02", "Analyse und Erkenntnisse")
    create_chart_placeholder_slide(prs, "Datenanalyse")
    create_key_insight_slide(
        prs,
        "Verhaltensökonomische Interventionen können die Conversion Rate um 23% steigern",
        [
            "Basierend auf 15 A/B-Tests mit >10'000 Teilnehmenden",
            "Grösster Effekt durch optimierte Default-Optionen",
            "ROI von 4.2x innerhalb von 6 Monaten",
        ]
    )
    create_section_slide(prs, "03", "Empfehlungen")
    create_content_slide(
        prs,
        "Handlungsempfehlungen",
        [
            "Implementierung eines Nudge-Frameworks für alle Kundenkontaktpunkte",
            "Schulung der Teams in verhaltensökonomischen Grundlagen",
            "Aufbau eines systematischen A/B-Testing Prozesses",
            "Quartalsweise Review und Optimierung der Massnahmen",
        ]
    )
    create_quote_slide(
        prs,
        "Der beste Weg, Verhalten zu ändern, ist nicht, Menschen zu überzeugen, sondern es ihnen leicht zu machen.",
        "Richard Thaler"
    )
    create_section_slide(prs, "04", "Nächste Schritte")
    create_content_slide(
        prs,
        "Projektplan",
        [
            "Phase 1: Discovery Workshop (2 Wochen)",
            "Phase 2: Analyse und Konzeption (4 Wochen)",
            "Phase 3: Pilotierung (6 Wochen)",
            "Phase 4: Rollout und Skalierung (8 Wochen)",
        ]
    )
    create_contact_slide(prs)

    # Save
    prs.save(output_path)
    print(f"✅ Template erstellt: {output_path}")
    print(f"   → {len(prs.slides)} Folien")

    return output_path


if __name__ == "__main__":
    import sys

    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_file = os.path.join(output_dir, "FehrAdvice_Template.pptx")

    if len(sys.argv) > 1:
        output_file = sys.argv[1]

    create_fehradvice_template(output_file)
